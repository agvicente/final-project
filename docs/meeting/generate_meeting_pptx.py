#!/usr/bin/env python3
"""
Gera apresentacao PowerPoint para reuniao com orientador.
Uso: python docs/meeting/generate_meeting_pptx.py

Estrutura (22 slides):
 1. Capa
 2. Objetivo da reuniao
 3. Agenda
 4. Pipeline Streaming
 5. Fundamentacao Teorica - TEDA (Angelov 2014)
 6. MicroTEDAclus - Algoritmo Principal (Maia 2020)
 7. Metodologia - Ablation Study
 8. C01 Baseline (+plot)
 9. C02 3 Hipoteses
10. C02-S3 Janelas Temporais (+plot +citacoes)
11. Framework Dimensoes Anomalas 1/3 - Definicao e Hipotese
12. Framework Dimensoes Anomalas 2/3 - Metodologia e Aplicacao
13. Framework Dimensoes Anomalas 3/3 - Validacao Estatistica Proposta
14. C03-S4 Features Comportamentais (+plot)
15. C04 5 Adaptacoes Tecnicas (+citacoes)
16. Por que a Formula Original Falha (demonstracao numerica)
17. C04 Impacto Quantitativo (+plot)
18. Consolidado - Melhor Resultado por Ataque
19. Contribuicoes Cientificas
20. O que Aprendemos - 6 Insights
21. Proximos Passos - Decisao Necessaria
22. Referencias Principais
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Paths
SCRIPT_DIR = Path(__file__).parent
PROJECT_ROOT = SCRIPT_DIR.parent.parent
OUTPUT_PATH = SCRIPT_DIR / "2026-03-19-advisor-meeting.pptx"

PLOTS_C01 = PROJECT_ROOT / "experiments" / "results" / "campaign-01" / "plots"
PLOTS_C02 = PROJECT_ROOT / "experiments" / "results" / "campaign-02" / "plots"
PLOTS_C03 = PROJECT_ROOT / "experiments" / "results" / "campaign-03" / "plots"
PLOTS_C04 = PROJECT_ROOT / "experiments" / "results" / "campaign-04" / "plots"

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MEDIUM_BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
LIGHT_GRAY = RGBColor(0x90, 0x90, 0x90)
GREEN = RGBColor(0x27, 0xAE, 0x60)
DARK_GREEN = RGBColor(0x1B, 0x5E, 0x20)
RED = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
GOLD = RGBColor(0xB8, 0x86, 0x0B)

# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ============================================================
# HELPER FUNCTIONS
# ============================================================

def set_slide_bg(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", italic=False):
    """Add a text box to slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=18, bold=False, color=BLACK,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6), italic=False):
    """Add paragraph to existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_table(slide, left, top, width, height, rows, cols):
    """Add table to slide."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table


def style_header_row(table, color=DARK_BLUE, font_color=WHITE):
    """Style the header row of a table."""
    for i, cell in enumerate(table.rows[0].cells):
        cell.fill.solid()
        cell.fill.fore_color.rgb = color
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = font_color
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.name = "Calibri"


def style_data_cell(cell, font_size=13, bold=False, color=BLACK,
                    alignment=PP_ALIGN.CENTER, italic=False):
    """Style a data cell."""
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def set_cell(table, row, col, text, font_size=13, bold=False, color=BLACK,
             italic=False, alignment=PP_ALIGN.CENTER):
    """Set cell text with styling."""
    cell = table.cell(row, col)
    cell.text = text
    style_data_cell(cell, font_size=font_size, bold=bold, color=color,
                    italic=italic, alignment=alignment)


def add_title_bar(slide, title_text):
    """Add a colored title bar at the top."""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.9),
                 title_text, font_size=30, bold=True, color=WHITE)


def add_citation_footer(slide, citation_text):
    """Add small citation text at the bottom of a slide."""
    add_text_box(slide, Inches(0.5), Inches(7.05), Inches(12.5), Inches(0.4),
                 citation_text, font_size=10, color=LIGHT_GRAY, italic=True)


def add_insight_box(slide, left, top, width, height, title, body, color=MEDIUM_BLUE):
    """Add a highlighted insight callout box."""
    shape = slide.shapes.add_shape(5, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)

    tf = add_text_box(slide, left + Inches(0.2), top + Inches(0.1),
                      width - Inches(0.4), height - Inches(0.2),
                      title, font_size=14, bold=True, color=color)
    add_paragraph(tf, body, font_size=12, color=DARK_GRAY)


# ============================================================
# PRESENTATION BUILDER
# ============================================================

def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]

    # ==================================================================
    # SLIDE 1: CAPA
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BLUE)

    add_text_box(slide, Inches(1), Inches(1.3), Inches(11), Inches(1.5),
                 "Progresso Experimental — Fase 2",
                 font_size=42, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(2.7), Inches(11), Inches(1.5),
                 "Deteccao de Intrusao em IoT com\nClustering Evolutivo em Streaming",
                 font_size=26, color=LIGHT_BLUE,
                 alignment=PP_ALIGN.CENTER)

    # Badge de volume experimental
    shape = slide.shapes.add_shape(5, Inches(3.5), Inches(4.8),
                                   Inches(6.3), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MEDIUM_BLUE
    shape.line.fill.background()
    add_text_box(slide, Inches(3.5), Inches(4.9), Inches(6.3), Inches(0.7),
                 "167 experimentos  |  4 campanhas  |  6 steps de ablation",
                 font_size=18, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    tf = add_text_box(slide, Inches(1), Inches(5.9), Inches(11), Inches(1.0),
                      "Augusto Custodio Vicente",
                      font_size=18, color=WHITE,
                      alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, "Mestrado PPGEE / UFMG",
                  font_size=16, color=LIGHT_BLUE,
                  alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, "Orientadores: Prof. Frederico Gadelha Guimaraes | Profa. Renata Lopes Rosa",
                  font_size=13, color=LIGHT_BLUE,
                  alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, "Reuniao de Orientacao",
                  font_size=13, color=LIGHT_BLUE,
                  alignment=PP_ALIGN.CENTER)

    # ==================================================================
    # SLIDE 2: OBJETIVO DA REUNIAO (NOVO)
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Objetivo desta Reuniao")

    # Caixa de contexto
    shape = slide.shapes.add_shape(5, Inches(0.5), Inches(1.5),
                                   Inches(12.3), Inches(1.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    shape.line.fill.background()

    tf = add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.8), Inches(1.2),
                      "Contexto: ~5 semanas ate a defesa (prazo maio 2026)",
                      font_size=18, bold=True, color=DARK_BLUE)
    add_paragraph(tf,
        "167 experimentos executados em 4 campanhas revelaram limites fundamentais da deteccao por anomalia, "
        "alinhados com literatura canonica, e identificaram 5 adaptacoes tecnicas originais.",
        font_size=14, color=DARK_GRAY)

    # Tres pontos em colunas
    add_text_box(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
                 "Preciso desta reuniao para:", font_size=20, bold=True, color=DARK_BLUE)

    objectives = [
        ("1.", "Validar achados",
         "Confirmar que o diagnostico do problema (representacao per-flow insuficiente) "
         "e as 5 adaptacoes tecnicas no MicroTEDAclus constituem contribuicao da dissertacao.",
         MEDIUM_BLUE),
        ("2.", "Discutir enquadramento",
         "Como apresentar os resultados negativos como contribuicao cientifica positiva, "
         "seguindo Sommer & Paxson (2010) e Arp et al. (2022).",
         GREEN),
        ("3.", "Decidir proximos passos",
         "Opcao A: S5 (Two-Stage Detection) + escrita. "
         "Opcao B: Consolidar + escrever. "
         "Qual faz mais sentido dado o prazo?",
         ORANGE),
    ]
    for i, (num, title, desc, color) in enumerate(objectives):
        x = Inches(0.5) + Inches(i * 4.3)
        y = Inches(3.9)

        shape = slide.shapes.add_shape(5, x, y, Inches(4.1), Inches(2.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        tf = add_text_box(slide, x + Inches(0.2), y + Inches(0.1),
                          Inches(3.8), Inches(0.5),
                          num, font_size=28, bold=True, color=color)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.8),
                     Inches(3.8), Inches(0.5),
                     title, font_size=18, bold=True, color=DARK_BLUE)
        add_text_box(slide, x + Inches(0.2), y + Inches(1.4),
                     Inches(3.8), Inches(1.5),
                     desc, font_size=12, color=DARK_GRAY)

    add_citation_footer(slide,
        'Sommer & Paxson (2010) "Outside the Closed World" IEEE S&P  |  '
        'Arp et al. (2022) "Dos and Don\'ts of ML in Computer Security" USENIX Security')

    # ==================================================================
    # SLIDE 3: AGENDA
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Agenda")

    items = [
        ("1.", "Pipeline Streaming", "PCAP -> Kafka -> MicroTEDAclus"),
        ("2.", "Fundamentacao Teorica", "TEDA (Angelov 2014) + MicroTEDAclus (Maia 2020)"),
        ("3.", "Metodologia", "Ablation study cumulativo (6 steps)"),
        ("4.", "Resultados por Campanha", "C01, C02 (S1/S2/S3), C03 (S4), C04"),
        ("5.", "Framework: Dimensoes Anomalas", "Hierarquia de detectabilidade"),
        ("6.", "Contribuicao Tecnica (C04)", "5 adaptacoes do MicroTEDAclus"),
        ("7.", "Contribuicoes Cientificas", "4 contribuicoes consolidadas"),
        ("8.", "Proximos Passos", "Decisao: S5 ou consolidar + escrever"),
    ]
    for i, (num, title, desc) in enumerate(items):
        y = Inches(1.6) + Inches(i * 0.65)
        add_text_box(slide, Inches(1.0), y, Inches(0.5), Inches(0.5),
                     num, font_size=20, bold=True, color=MEDIUM_BLUE)
        tf = add_text_box(slide, Inches(1.5), y, Inches(10), Inches(0.5),
                          title, font_size=20, bold=True, color=DARK_GRAY)
        add_paragraph(tf, desc, font_size=14, color=LIGHT_GRAY)

    # ==================================================================
    # SLIDE 4: PIPELINE STREAMING
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Pipeline Streaming")

    # 6 stages, evenly spaced across slide width (13.333")
    # box width = 1.75", stride = 2.1", start = 0.42"
    # last box ends at 0.42 + 5*2.1 + 1.75 = 12.67 (fits)
    box_width = Inches(1.75)
    box_height = Inches(1.2)
    stride = Inches(2.1)
    start_x = Inches(0.42)
    box_y = Inches(2.0)

    stage_labels = [
        "PCAPs\nCICIoT2023",
        "Kafka\nProducer",
        "Kafka\nBroker",
        "Flow\nConsumer",
        "MicroTEDA\nclus",
        "Metricas\nPrequential",
    ]
    stage_x_positions = []
    for i, label in enumerate(stage_labels):
        x = start_x + stride * i
        stage_x_positions.append(x)
        shape = slide.shapes.add_shape(5, x, box_y, box_width, box_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = MEDIUM_BLUE
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

    # Arrows between stages (Unicode right arrow)
    for i in range(len(stage_labels) - 1):
        x_start = stage_x_positions[i] + box_width
        x_end = stage_x_positions[i+1]
        x_mid = (x_start + x_end) / 2
        add_text_box(slide, Emu(int(x_mid - Inches(0.18))), Inches(2.32),
                     Inches(0.36), Inches(0.56), "\u2192",
                     font_size=32, bold=True, color=DARK_BLUE,
                     alignment=PP_ALIGN.CENTER)

    # 6 details, one per stage, aligned horizontally with each box
    details = [
        ("5 tipos de ataque", "Benigno + DDoS\n(ICMP/SYN/TCP)\n+ Mirai + Recon"),
        ("Leitura de PCAP", "Envio em ordem\nde timestamp"),
        ("Topico 'packets'", "Ordem preservada\n(reprodutibilidade)"),
        ("Reconstrucao de flows", "17 features\nTimeout 60s"),
        ("Clustering evolutivo", "Tipicalidade +\nChebyshev dinamico\n(Maia 2020)"),
        ("Test-then-train", "Ground truth IP\nFading α=0.01\n(Gama 2013)"),
    ]
    for i, (title, desc) in enumerate(details):
        x = stage_x_positions[i]
        tf = add_text_box(slide, x, Inches(3.6), box_width, Inches(2.0),
                          title, font_size=12, bold=True, color=DARK_BLUE,
                          alignment=PP_ALIGN.CENTER)
        add_paragraph(tf, desc, font_size=10, color=DARK_GRAY,
                      alignment=PP_ALIGN.CENTER)

    add_insight_box(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.85),
        "Decisoes de design",
        "Kafka garante reprodutibilidade (ordem de insercao) + desacoplamento producer/consumer/detector. "
        "Prequential evita data leakage e respeita dependencias temporais do stream.")

    add_citation_footer(slide,
        'Neto et al. (2023) "CICIoT2023" Sensors 23(13)  |  '
        'Gama et al. (2013) "On Evaluating Stream Learning Algorithms" Machine Learning 90(3)')

    # ==================================================================
    # SLIDE 5: FUNDAMENTACAO TEORICA - TEDA (NOVO)
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Fundamentacao Teorica — TEDA (Angelov 2014)")

    # Equacao principal da eccentricidade
    shape = slide.shapes.add_shape(5, Inches(0.5), Inches(1.5),
                                   Inches(12.3), Inches(1.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    shape.line.color.rgb = DARK_BLUE
    shape.line.width = Pt(1.5)

    tf = add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.5),
                      "Eccentricidade (quao diferente um ponto e):",
                      font_size=14, bold=True, color=DARK_BLUE)
    add_text_box(slide, Inches(0.8), Inches(2.05), Inches(11.7), Inches(0.7),
                 "ξ(x_k) = 1/k + ||x_k - μ_k||² / (k × σ²_k)",
                 font_size=22, bold=True, color=BLACK,
                 font_name="Cambria Math", alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(2.65), Inches(11.7), Inches(0.4),
                 "Tipicalidade:  τ(x_k) = 1 − ξ(x_k)    |    "
                 "Rejeicao via Chebyshev:  ζ(x) ≤ (m² + 1) / (2n)",
                 font_size=13, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # Tres caracteristicas fundamentais
    characteristics = [
        ("Nao-parametrico",
         "Nao assume gaussianidade, independencia, ou amostras infinitas. "
         "Usa apenas estrutura espacial mutua dos dados.",
         MEDIUM_BLUE),
        ("Single-pass O(1)",
         "Atualizacao recursiva via Welford (1962): mantem apenas (n, μ, σ²). "
         "Adequado para streaming de alta vazao.",
         GREEN),
        ("Chebyshev sem distribuicao",
         "P(|X-μ| ≥ mσ) ≤ 1/m² vale para QUALQUER distribuicao. "
         "Com m=3, 89% dos pontos ficam dentro (Amidan 2005).",
         ORANGE),
    ]
    for i, (title, desc, color) in enumerate(characteristics):
        x = Inches(0.5) + Inches(i * 4.3)
        y = Inches(3.4)

        shape = slide.shapes.add_shape(5, x, y, Inches(4.1), Inches(2.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = color
        shape.line.width = Pt(1.5)

        tf = add_text_box(slide, x + Inches(0.2), y + Inches(0.15),
                          Inches(3.7), Inches(0.5),
                          title, font_size=15, bold=True, color=color)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.7),
                     Inches(3.7), Inches(1.3),
                     desc, font_size=12, color=DARK_GRAY)

    # Ponte narrativa para o proximo slide (MicroTEDAclus)
    shape = slide.shapes.add_shape(5, Inches(0.5), Inches(5.6),
                                   Inches(12.3), Inches(1.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
    shape.line.color.rgb = RED
    shape.line.width = Pt(1.5)

    tf = add_text_box(slide, Inches(0.8), Inches(5.7), Inches(11.8), Inches(0.5),
                      "Limitacao do TEDA basico — centro global unico",
                      font_size=14, bold=True, color=RED)
    add_paragraph(tf,
        "Um unico μ e σ² para toda a populacao. Grupos distintos parecem igualmente "
        "excentricos; outliers contaminam as estatisticas globais; nao ha mecanismo "
        "para separar padroes de ataque distintos.",
        font_size=12, color=DARK_GRAY)
    add_paragraph(tf,
        "Solucao: MicroTEDAclus (proximo slide)",
        font_size=12, bold=True, italic=True, color=DARK_BLUE)

    add_citation_footer(slide,
        'Angelov (2014) "Outside the Box" JAMRIS 8(2)  |  '
        'Welford (1962) Technometrics 4(3)  |  '
        'Amidan et al. (2005) "Data Outlier Detection Using Chebyshev" IEEE Aerospace')

    # ==================================================================
    # SLIDE 6: MICROTEDACLUS — ALGORITMO PRINCIPAL (NOVO)
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "MicroTEDAclus (Maia 2020) — Algoritmo Principal")

    # Intro
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
                 "Extensao do TEDA: N micro-clusters com estatisticas ISOLADAS (μ_i, σ²_i, n_i)",
                 font_size=16, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(1.75), Inches(12.3), Inches(0.4),
                 "Cada cluster protege suas proprias estatisticas — outliers criam novos clusters ao inves de contaminar os existentes",
                 font_size=12, italic=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # ---------- LEFT COLUMN: Fluxo de decisao ----------
    add_text_box(slide, Inches(0.5), Inches(2.35), Inches(6.0), Inches(0.4),
                 "Fluxo de decisao por ponto x:",
                 font_size=14, bold=True, color=MEDIUM_BLUE)

    # Step 1: Ponto chega
    shape = slide.shapes.add_shape(5, Inches(0.5), Inches(2.85),
                                   Inches(6.0), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    shape.line.color.rgb = MEDIUM_BLUE
    shape.line.width = Pt(1)
    tf = add_text_box(slide, Inches(0.6), Inches(2.9),
                      Inches(5.8), Inches(0.4),
                      "1. Ponto x chega (vetor de 17 features)",
                      font_size=12, bold=True, color=DARK_BLUE)

    # Arrow
    add_text_box(slide, Inches(3.3), Inches(3.38), Inches(0.4), Inches(0.3),
                 "\u2193", font_size=16, bold=True, color=MEDIUM_BLUE,
                 alignment=PP_ALIGN.CENTER)

    # Step 2: Chebyshev test
    shape = slide.shapes.add_shape(5, Inches(0.5), Inches(3.65),
                                   Inches(6.0), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    shape.line.color.rgb = MEDIUM_BLUE
    shape.line.width = Pt(1)
    tf = add_text_box(slide, Inches(0.6), Inches(3.7),
                      Inches(5.8), Inches(0.4),
                      "2. Para cada micro-cluster MC_i:",
                      font_size=12, bold=True, color=DARK_BLUE)
    add_paragraph(tf, "   Calcula ζ_i(x) = ξ_i(x) / 2",
                  font_size=11, color=DARK_GRAY)
    add_paragraph(tf, "   Aceita se  ζ_i ≤ (m_i(k)² + 1) / (2·n_i)",
                  font_size=11, color=DARK_GRAY)

    # Arrow
    add_text_box(slide, Inches(3.3), Inches(4.6), Inches(0.4), Inches(0.3),
                 "\u2193", font_size=16, bold=True, color=MEDIUM_BLUE,
                 alignment=PP_ALIGN.CENTER)

    # Step 3 (bifurcation: SIM / NAO)
    # Left: SIM branch
    shape = slide.shapes.add_shape(5, Inches(0.5), Inches(4.9),
                                   Inches(2.9), Inches(1.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(1.5)
    tf = add_text_box(slide, Inches(0.6), Inches(4.95),
                      Inches(2.7), Inches(0.95),
                      "3a. Algum MC aceita:", font_size=11, bold=True, color=GREEN)
    add_paragraph(tf, "Atualiza SO o melhor",
                  font_size=10, color=DARK_GRAY)
    add_paragraph(tf, "(maior typicality)",
                  font_size=10, color=DARK_GRAY)
    add_paragraph(tf, "\u2192 NORMAL",
                  font_size=11, bold=True, color=DARK_GREEN)

    # Right: NAO branch
    shape = slide.shapes.add_shape(5, Inches(3.6), Inches(4.9),
                                   Inches(2.9), Inches(1.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
    shape.line.color.rgb = RED
    shape.line.width = Pt(1.5)
    tf = add_text_box(slide, Inches(3.7), Inches(4.95),
                      Inches(2.7), Inches(0.95),
                      "3b. Nenhum MC aceita:", font_size=11, bold=True, color=RED)
    add_paragraph(tf, "Cria novo MC",
                  font_size=10, color=DARK_GRAY)
    add_paragraph(tf, "(centro = x)",
                  font_size=10, color=DARK_GRAY)
    add_paragraph(tf, "\u2192 ANOMALIA",
                  font_size=11, bold=True, color=RED)

    # ---------- RIGHT COLUMN: vs TEDA comparison table ----------
    add_text_box(slide, Inches(6.8), Inches(2.35), Inches(6.2), Inches(0.4),
                 "MicroTEDAclus vs TEDA basico:",
                 font_size=14, bold=True, color=MEDIUM_BLUE)

    table = add_table(slide, Inches(6.8), Inches(2.85),
                      Inches(6.3), Inches(3.1), 5, 3)
    for i, h in enumerate(["Aspecto", "TEDA basico", "MicroTEDAclus"]):
        table.cell(0, i).text = h
    style_header_row(table)

    comparison = [
        ("Centro", "1 global (μ, σ²)", "N independentes"),
        ("Outliers", "Contaminam\nestatisticas", "Criam novos\nclusters (isolados)"),
        ("Numero K", "Implicito (= 1)", "Automatico,\ncresce dinamicamente"),
        ("Concept drift", "Limitado", "Nativo\n(novos padroes = novos MCs)"),
    ]
    for r, (aspect, teda, mtc) in enumerate(comparison):
        set_cell(table, r+1, 0, aspect, font_size=11, bold=True, color=DARK_BLUE)
        set_cell(table, r+1, 1, teda, font_size=10, color=RED)
        set_cell(table, r+1, 2, mtc, font_size=10, color=DARK_GREEN)

    # ---------- BOTTOM: Threshold dinamico box ----------
    shape = slide.shapes.add_shape(5, Inches(0.5), Inches(6.1),
                                   Inches(12.5), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE1)
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(1.5)

    tf = add_text_box(slide, Inches(0.7), Inches(6.15), Inches(12.1), Inches(0.4),
                      "Threshold dinamico m(k) — protege clusters jovens, exige de clusters maduros",
                      font_size=13, bold=True, color=GOLD)
    add_paragraph(tf,
        "m(k) = 3 / (1 + e^(-0.007 × (k−100)))     "
        "|     k=1: m≈0.6 (permissivo)     |     k=100: m=1.5 (transicao)     "
        "|     k\u2192\u221e: m\u21923 (89% Chebyshev, estrito)",
        font_size=11, color=DARK_GRAY)

    add_citation_footer(slide,
        'Maia et al. (2020) "Evolving Clustering Based on Mixture of Typicalities" FGCS 106  |  '
        'Kohonen (1990) "Self-Organizing Map" Proc IEEE 78(9) - fundamento do update seletivo  |  '
        'Cao et al. (2006) "DenStream" SDM')

    # ==================================================================
    # SLIDE 7: METODOLOGIA — ABLATION STUDY
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Metodologia — Ablation Study Cumulativo")

    table = add_table(slide, Inches(0.3), Inches(1.4),
                      Inches(12.7), Inches(4.2), 7, 5)

    headers = ["Step", "Variavel Testada", "Opcoes", "Resultado", "Decisao"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
    style_header_row(table)

    data = [
        ["C01", "Algoritmo", "TEDA vs MicroTEDAclus", "MicroTEDAclus 26x melhor", "MicroTEDAclus"],
        ["C02-S1", "Ground Truth", "Phase vs IP", "ICMP: 4%->27% (L4 corrigido)", "IP GT"],
        ["C02-S2", "Features/flow", "v1(17) vs v2(25) vs v3(32)", "ZERO impacto (±1pp)", "v1 (Occam)"],
        ["C02-S3", "Granularidade", "Per-flow vs Window", "Recall 10-20x melhor", "Window"],
        ["C03-S4", "Features/janela", "v1(12) vs v2(19 comport.)", "Misto: 2/5 melhor, 2/5 pior", "Em analise"],
        ["C04", "Implementacao", "Propria vs Original (Maia)", "Original FPR 42-75%", "Propria superior"],
    ]
    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            set_cell(table, r+1, c, val, font_size=12)

    add_insight_box(slide, Inches(0.3), Inches(5.8), Inches(12.7), Inches(1.1),
        "Principio metodologico",
        "Cada step altera UMA unica variavel e congela a melhor decisao para o proximo "
        "-> causalidade isolada. Total: 167 runs. Arp et al. (2022) chamam isso de rigor essencial "
        "para ML em seguranca, evitando resultados inflados.")

    add_citation_footer(slide,
        'Arp et al. (2022) "Dos and Don\'ts of ML in Computer Security" USENIX Security 2022  |  '
        'Pendlebury et al. (2019) "TESSERACT" USENIX Security 2019')

    # ==================================================================
    # SLIDE 7: C01 BASELINE (com plot)
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Campaign-01 — Baseline (17 runs)")

    # Plot: anomaly rate invariante
    plot_path = PLOTS_C01 / "01_anomaly_rate_invariant.png"
    if plot_path.exists():
        slide.shapes.add_picture(
            str(plot_path),
            Inches(0.3), Inches(1.4),
            Inches(7.5), Inches(5.4)
        )
    else:
        add_text_box(slide, Inches(1), Inches(3), Inches(6), Inches(1),
                     f"[Plot nao encontrado: {plot_path.name}]",
                     font_size=16, color=RED)

    # Side: key finding
    shape = slide.shapes.add_shape(5, Inches(8.1), Inches(1.4),
                                   Inches(5.0), Inches(1.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
    shape.line.fill.background()
    tf = add_text_box(slide, Inches(8.3), Inches(1.5), Inches(4.8), Inches(1.2),
                      "Achado principal:",
                      font_size=14, bold=True, color=RED)
    add_paragraph(tf, "Anomaly rate invariante (~3.5%) com OU sem ataque",
                  font_size=13, color=DARK_GRAY)

    # Results table
    table = add_table(slide, Inches(8.1), Inches(3.0),
                      Inches(5.0), Inches(2.0), 4, 3)
    for i, h in enumerate(["Cenario", "Resultado", "Alvo"]):
        table.cell(0, i).text = h
    style_header_row(table)

    c01_data = [
        ["Benigno FPR", "3.5%", "<= 5%"],
        ["Ataques Recall", "~3-4%", ">= 80%"],
        ["TEDA vs MTC", "26x", "—"],
    ]
    for r, row_data in enumerate(c01_data):
        for c, val in enumerate(row_data):
            color = GREEN if val == "3.5%" else (RED if "3-4%" in val else BLACK)
            bold = color != BLACK
            set_cell(table, r+1, c, val, font_size=12, bold=bold, color=color)

    # Diagnostico
    tf = add_text_box(slide, Inches(8.1), Inches(5.2), Inches(5.0), Inches(1.7),
                      "Diagnostico:",
                      font_size=14, bold=True, color=DARK_BLUE)
    add_paragraph(tf,
        'O detector encontra outliers estatisticos, nao ataques. '
        'Sommer & Paxson (2010) previram exatamente isso: '
        '"gap semantico entre anomalia e malicia".',
        font_size=12, color=DARK_GRAY)

    add_citation_footer(slide,
        'Sommer & Paxson (2010) "Outside the Closed World: On Using ML for Network IDS" IEEE S&P 2010  |  '
        'Axelsson (2000) "The Base-Rate Fallacy" ACM TISSEC 3(3)')

    # ==================================================================
    # SLIDE 8: C02 — 3 HIPOTESES
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Campaign-02 — 3 Hipoteses (72 runs)")

    table = add_table(slide, Inches(0.3), Inches(1.4),
                      Inches(12.7), Inches(4.5), 4, 5)

    for i, h in enumerate(["Step", "Variavel", "Resultado", "Decisao", "Impacto"]):
        table.cell(0, i).text = h
    style_header_row(table)

    s_data = [
        ["S1 — GT por IP", "phase -> IP",
         "ICMP: 4%->27% Recall\nResto inalterado\nFPR estavel ~3-4%",
         "Adotar IP GT",
         "Corrige medicao,\nmas nao resolve\ndeteccao"],
        ["S2 — Features", "v1(17)->v2(25)->v3(32)",
         "ZERO impacto (±1pp)\nem TODOS os ataques",
         "Manter v1\n(Occam's razor)",
         "Features per-flow\nsaturaram:\nproblema estrutural"],
        ["S3 — Janelas", "Per-flow -> Window\n(5s, 10s, 30s, 60s)",
         "SYN: 3%->54% (15x)\nRecon: 4%->45% (10x)\nMAS FPR explode (58% @60s)",
         "Direcao certa,\nfeatures insuficientes",
         "Muda a pergunta\nfundamental"],
    ]
    for r, row_data in enumerate(s_data):
        for c, val in enumerate(row_data):
            set_cell(table, r+1, c, val, font_size=11)

    add_insight_box(slide, Inches(0.3), Inches(6.05), Inches(12.7), Inches(0.95),
        "Insight central de S3",
        'Mudanca de granularidade muda a PERGUNTA do detector: '
        '"este flow e anomalo?" (resposta: NAO, flows DDoS parecem normais) -> '
        '"este IP comporta-se de forma anomala?" (resposta: SIM, padrao de flood emerge).')

    add_citation_footer(slide,
        'Lakhina, Crovella & Diot (2004/2005) ACM SIGCOMM: seminal em analise agregada de anomalias de rede')

    # ==================================================================
    # SLIDE 9: C02-S3 JANELAS (com plot)
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "C02-S3 — Recall por Janela Temporal")

    plot_path = PLOTS_C02 / "04_s3_recall_by_window.png"
    if plot_path.exists():
        slide.shapes.add_picture(
            str(plot_path),
            Inches(0.3), Inches(1.4),
            Inches(7.5), Inches(5.2)
        )

    tf = add_text_box(slide, Inches(8.2), Inches(1.5), Inches(4.8), Inches(0.6),
                      "Melhorias dramaticas de Recall:",
                      font_size=16, bold=True, color=DARK_BLUE)

    improvements = [
        ("DDoS-SYN:", "3.5% -> 53.9% (15x)", GREEN),
        ("Recon:", "4.5% -> 45.3% (10x)", GREEN),
        ("Mirai:", "1.7% -> 33.3% (20x)", GREEN),
    ]
    for label, val, color in improvements:
        add_paragraph(tf, f"  {label} {val}", font_size=14, bold=True, color=color)

    add_paragraph(tf, "", font_size=6)
    add_paragraph(tf, "Problema:", font_size=14, bold=True, color=RED)
    add_paragraph(tf, "  FPR explode com janela grande", font_size=13, color=RED)
    add_paragraph(tf, "  (58% @60s benigno)", font_size=13, color=RED)

    add_paragraph(tf, "", font_size=6)
    add_paragraph(tf, "Suporte na literatura:", font_size=14, bold=True, color=DARK_BLUE)
    add_paragraph(tf,
        "Li et al. (2023) e Goldschmidt & Kucera (2024) "
        "reportam ganhos similares com window-based features para deteccao de DDoS.",
        font_size=11, color=DARK_GRAY, italic=True)

    add_citation_footer(slide,
        'Li et al. (2023) "Window-based Feature Extraction for DDoS" Sci China Info Sci 66  |  '
        'Goldschmidt & Kucera (2024) "Windower" NOMS 2024  |  '
        'Lakhina et al. (2005) "Mining Anomalies Using Traffic Feature Distributions" SIGCOMM')

    # ==================================================================
    # SLIDE 11: FRAMEWORK DIMENSOES ANOMALAS 1/3 — DEFINICAO E HIPOTESE
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, 'Framework Proposto: Dimensoes Anomalas (1/3)')

    # CAVEAT no topo
    shape = slide.shapes.add_shape(5, Inches(0.5), Inches(1.3),
                                   Inches(12.3), Inches(0.75))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE1)
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(1.5)
    tf = add_text_box(slide, Inches(0.7), Inches(1.35), Inches(12.0), Inches(0.65),
                      "Sintese analitica proposta neste trabalho — NAO e conceito estabelecido na literatura",
                      font_size=13, bold=True, color=GOLD)
    add_paragraph(tf,
        "Ferramenta construida para organizar a hierarquia de detectabilidade "
        "observada nos 167 experimentos",
        font_size=11, italic=True, color=DARK_GRAY)

    # DEFINICAO — coluna esquerda
    add_text_box(slide, Inches(0.5), Inches(2.25), Inches(6.2), Inches(0.45),
                 "Definicao: o que e uma 'dimensao anomala'?",
                 font_size=15, bold=True, color=MEDIUM_BLUE)

    tf = add_text_box(slide, Inches(0.5), Inches(2.75), Inches(6.2), Inches(2.6),
                      "Uma feature f_j e 'anomala' para o ataque A quando:",
                      font_size=12, color=DARK_GRAY)
    add_paragraph(tf,
        "   a distribuicao de f_j durante o ataque A difere substancialmente "
        "da distribuicao de f_j em trafego benigno.",
        font_size=12, color=DARK_GRAY)
    add_paragraph(tf, "", font_size=4)
    add_paragraph(tf, "Formal:",
                  font_size=12, bold=True, color=DARK_BLUE)
    add_paragraph(tf,
        "   D_benign(f_j)  ≠  D_attack(f_j)     (qualitativamente distinguiveis)",
        font_size=12, color=BLACK)
    add_paragraph(tf, "", font_size=4)
    add_paragraph(tf, "Contagem por ataque:",
                  font_size=12, bold=True, color=DARK_BLUE)
    add_paragraph(tf,
        "   d_A = numero de features f_j cujas distribuicoes benigna e "
        "de ataque A sao visualmente separaveis",
        font_size=12, color=DARK_GRAY)

    # EXEMPLO CONCRETO — coluna direita
    add_text_box(slide, Inches(6.9), Inches(2.25), Inches(6.0), Inches(0.45),
                 "Exemplo concreto — feature dst_port_entropy:",
                 font_size=15, bold=True, color=MEDIUM_BLUE)

    # Caixa do exemplo
    shape = slide.shapes.add_shape(5, Inches(6.9), Inches(2.75),
                                   Inches(6.0), Inches(2.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    shape.line.color.rgb = MEDIUM_BLUE
    shape.line.width = Pt(1)

    tf = add_text_box(slide, Inches(7.1), Inches(2.85), Inches(5.7), Inches(2.4),
                      "Trafego benigno IoT:", font_size=12, bold=True, color=GREEN)
    add_paragraph(tf, "   ~1.2 bits em media (2-3 portas ativas: DNS, HTTPS, MQTT)",
                  font_size=11, color=DARK_GRAY)
    add_paragraph(tf, "", font_size=4)
    add_paragraph(tf, "Durante Recon-PortScan:",
                  font_size=12, bold=True, color=RED)
    add_paragraph(tf, "   ~6.5 bits (100+ portas probadas sequencialmente)",
                  font_size=11, color=DARK_GRAY)
    add_paragraph(tf, "", font_size=4)
    add_paragraph(tf,
        "→ Distribuicoes sao claramente separaveis",
        font_size=12, bold=True, color=DARK_BLUE)
    add_paragraph(tf,
        "→ dst_port_entropy conta como 1 dimensao anomala para Recon",
        font_size=12, bold=True, color=DARK_BLUE)
    add_paragraph(tf, "", font_size=4)
    add_paragraph(tf,
        "Contra-exemplo: a mesma feature para DDoS-TCP flood tem "
        "distribuicao identica ao benigno (ambos usam poucas portas alvo) "
        "→ NAO conta.",
        font_size=10, italic=True, color=DARK_GRAY)

    # HIPOTESE CAUSAL — base inferior
    shape = slide.shapes.add_shape(5, Inches(0.5), Inches(5.55),
                                   Inches(12.4), Inches(1.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    shape.line.color.rgb = DARK_BLUE
    shape.line.width = Pt(1.5)

    tf = add_text_box(slide, Inches(0.7), Inches(5.6), Inches(12.0), Inches(1.4),
                      "Hipotese causal",
                      font_size=14, bold=True, color=DARK_BLUE)
    add_paragraph(tf,
        "Mais dimensoes anomalas  →  maior ||x - μ||²  →  maior ξ (eccentricidade)  "
        "→  maior probabilidade de rejeicao por Chebyshev  →  maior Recall",
        font_size=12, color=BLACK)
    add_paragraph(tf,
        "Intuicao: cada dimensao anomala contribui aditivamente para "
        "||x - μ||² = Σ(x_j - μ_j)², que entra direto na formula "
        "ξ = 1/k + ||x-μ||² / (k × σ²). Ataque com 0 dimensoes anomalas "
        "tem ||x - μ||² ≈ ||benigno - μ||² e fica preso no cluster benigno.",
        font_size=11, italic=True, color=DARK_GRAY)

    add_citation_footer(slide,
        'Inspiracao: Sommer & Paxson (2010) "Outside the Closed World" sobre gap semantico  |  '
        'Lakhina et al. (2005) "Mining Anomalies Using Traffic Feature Distributions" SIGCOMM')

    # ==================================================================
    # SLIDE 12: FRAMEWORK DIMENSOES ANOMALAS 2/3 — METODOLOGIA + APLICACAO
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Framework Proposto: Dimensoes Anomalas (2/3)")

    # Metodologia de contagem
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.4),
                 "Metodologia de contagem (analise qualitativa):",
                 font_size=14, bold=True, color=MEDIUM_BLUE)

    tf = add_text_box(slide, Inches(0.5), Inches(1.75), Inches(12.3), Inches(1.1),
                      "1. Inspecao visual das 19 features de janela (v2) para cada tipo de ataque "
                      "usando os dados agregados das campanhas C02-C03",
                      font_size=11, color=DARK_GRAY)
    add_paragraph(tf,
        "2. Cross-reference com literatura sobre caracteristicas de cada ataque "
        "(Zargar 2013 para DDoS; Nychis 2008 para scan; Antonakakis 2017 para Mirai)",
        font_size=11, color=DARK_GRAY)
    add_paragraph(tf,
        "3. Classificacao em 3 niveis: FORTE (distribuicoes claramente separaveis) / "
        "FRACA (sobreposicao parcial) / NULA (sobreposicao total). "
        "Apenas dimensoes FORTES entram na contagem d_A.",
        font_size=11, color=DARK_GRAY)
    add_paragraph(tf,
        "Limitacao: classificacao e qualitativa, baseada em inspecao + literatura. "
        "Nao ha teste estatistico formal (ver slide 3/3 para experimento proposto).",
        font_size=11, italic=True, color=RED)

    # Tabela hierarquica aplicada
    table = add_table(slide, Inches(0.3), Inches(3.35),
                      Inches(12.7), Inches(2.8), 6, 5)
    for i, h in enumerate(["Ataque", "d_A", "Features Fortes", "Recall Medido", "Referencia Literatura"]):
        table.cell(0, i).text = h
    style_header_row(table)

    hierarchy = [
        ("Recon-PortScan", "5+",
         "port_entropy↑, dst_div↑, flow_size↓, dur↓, fwd_only↑",
         "49%", "Nychis (2008)", GREEN),
        ("DDoS-ICMP (v2)", "3-4",
         "flows/s↑, payload_std↓, dst_ip_ent↓, small↑",
         "50%", "Wang (2022)", GREEN),
        ("Mirai-greeth", "2-3",
         "rate↑, target_concentration↑",
         "46%", "Antonakakis (2017)", ORANGE),
        ("DDoS-SYN", "1-2",
         "syn_ratio↑, unanswered↑",
         "38-62%", "Bellaiche (2012)", ORANGE),
        ("DDoS-TCP", "0",
         "TODAS sobrepoem com benigno",
         "0%", "Zargar (2013)", RED),
    ]
    for r, (attack, ndims, dims, recall, ref, color) in enumerate(hierarchy):
        set_cell(table, r+1, 0, attack, font_size=11, bold=True)
        set_cell(table, r+1, 1, ndims, font_size=15, bold=True, color=color)
        set_cell(table, r+1, 2, dims, font_size=9, alignment=PP_ALIGN.LEFT)
        set_cell(table, r+1, 3, recall, font_size=13, bold=True, color=color)
        set_cell(table, r+1, 4, ref, font_size=10, italic=True, color=DARK_GRAY)

    # Observacao
    add_insight_box(slide, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.8),
        "Padrao observado: Recall cresce monotonicamente com d_A  (5+→49%, 3-4→50%, 2-3→46%, 1-2→38-62%, 0→0%)",
        "Consistente com a hipotese causal, mas N=5 ataques nao permite validacao estatistica rigorosa "
        "-> motivacao para o experimento do proximo slide")

    add_citation_footer(slide,
        "Recall: medido nas campanhas C02-S3 e C03-S4  |  "
        "Features fortes: inspecao das ANALYSIS.md + literatura citada")

    # ==================================================================
    # SLIDE 13: FRAMEWORK DIMENSOES ANOMALAS 3/3 — VALIDACAO ESTATISTICA
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Framework Proposto: Dimensoes Anomalas (3/3)")

    # Intro
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
                 "Experimento proposto para validacao estatistica rigorosa (nao executado por restricao de tempo)",
                 font_size=14, italic=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

    # COLUNA ESQUERDA: Fase 1 — contagem estatistica via KS-test
    add_text_box(slide, Inches(0.5), Inches(2.0), Inches(6.2), Inches(0.4),
                 "Fase 1: Contagem estatistica via KS-test",
                 font_size=13, bold=True, color=MEDIUM_BLUE)

    shape = slide.shapes.add_shape(5, Inches(0.5), Inches(2.45),
                                   Inches(6.2), Inches(3.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    shape.line.color.rgb = MEDIUM_BLUE
    shape.line.width = Pt(1)

    tf = add_text_box(slide, Inches(0.7), Inches(2.55), Inches(5.8), Inches(3.4),
                      "Para cada feature f_j (j=1..19) e ataque A:",
                      font_size=11, bold=True, color=DARK_BLUE)
    add_paragraph(tf,
        "1. Extrair D_benign(f_j) e D_attack_A(f_j) dos dados",
        font_size=11, color=DARK_GRAY)
    add_paragraph(tf,
        "2. Aplicar Kolmogorov-Smirnov 2-sample test",
        font_size=11, color=DARK_GRAY)
    add_paragraph(tf,
        "   (nao-parametrico, distribuicao-livre)",
        font_size=10, italic=True, color=DARK_GRAY)
    add_paragraph(tf,
        "3. Registrar estatistica KS e p-valor",
        font_size=11, color=DARK_GRAY)
    add_paragraph(tf,
        "4. Contar f_j como 'anomala' se p < 0.01",
        font_size=11, color=DARK_GRAY)
    add_paragraph(tf,
        "   (correcao de Bonferroni para 19 testes)",
        font_size=10, italic=True, color=DARK_GRAY)
    add_paragraph(tf, "", font_size=4)
    add_paragraph(tf,
        "Resultado: d_A estatistico por ataque",
        font_size=11, bold=True, color=DARK_BLUE)
    add_paragraph(tf,
        "Substitui a contagem qualitativa por uma",
        font_size=10, italic=True, color=DARK_GRAY)
    add_paragraph(tf,
        "medicao objetiva e reproduzivel.",
        font_size=10, italic=True, color=DARK_GRAY)

    # COLUNA DIREITA: Fase 2 — teste de correlacao
    add_text_box(slide, Inches(6.9), Inches(2.0), Inches(6.2), Inches(0.4),
                 "Fase 2: Teste de correlacao d_A  ↔  Recall_A",
                 font_size=13, bold=True, color=MEDIUM_BLUE)

    shape = slide.shapes.add_shape(5, Inches(6.9), Inches(2.45),
                                   Inches(6.2), Inches(3.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(1)

    tf = add_text_box(slide, Inches(7.1), Inches(2.55), Inches(5.8), Inches(3.4),
                      "Com os d_A estatisticos da Fase 1:",
                      font_size=11, bold=True, color=DARK_GREEN)
    add_paragraph(tf,
        "1. Pares (d_A, Recall_A) para cada ataque A",
        font_size=11, color=DARK_GRAY)
    add_paragraph(tf,
        "2. Coeficiente de Spearman (nao-parametrico)",
        font_size=11, color=DARK_GRAY)
    add_paragraph(tf,
        "3. H0: nao ha correlacao entre d_A e Recall",
        font_size=11, color=DARK_GRAY)
    add_paragraph(tf,
        "   H1: correlacao positiva (rho > 0)",
        font_size=11, color=DARK_GRAY)
    add_paragraph(tf,
        "4. Reportar rho, p-valor, intervalo bootstrap",
        font_size=11, color=DARK_GRAY)
    add_paragraph(tf, "", font_size=4)
    add_paragraph(tf,
        "Limitacao: N=5 ataques",
        font_size=11, bold=True, color=RED)
    add_paragraph(tf,
        "Poder estatistico baixo, mas direcionalidade",
        font_size=10, italic=True, color=DARK_GRAY)
    add_paragraph(tf,
        "e reportavel. Extensao futura: ampliar para",
        font_size=10, italic=True, color=DARK_GRAY)
    add_paragraph(tf,
        "10+ ataques do CICIoT2023 (tempo adicional).",
        font_size=10, italic=True, color=DARK_GRAY)

    # BOTTOM: O que o experimento resolveria + viabilidade
    shape = slide.shapes.add_shape(5, Inches(0.5), Inches(6.2),
                                   Inches(12.5), Inches(0.85))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE1)
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(1.5)

    tf = add_text_box(slide, Inches(0.7), Inches(6.25), Inches(12.1), Inches(0.8),
                      "O que este experimento resolveria",
                      font_size=13, bold=True, color=GOLD)
    add_paragraph(tf,
        "(a) Substitui contagem qualitativa por estatistica objetiva.  "
        "(b) Valida ou falsifica a hipotese causal.  "
        "(c) Produz um 'indice de detectabilidade' por ataque reproduzivel.  "
        "Tempo estimado: ~2-3 dias (scripts de extracao ja existem nas campanhas anteriores).",
        font_size=10, color=DARK_GRAY)

    add_citation_footer(slide,
        'Massey (1951) "The Kolmogorov-Smirnov Test for Goodness of Fit" JASA  |  '
        'Spearman (1904) "The Proof and Measurement of Association Between Two Things" Am J Psych')

    # ==================================================================
    # SLIDE 11: C03-S4 FEATURES COMPORTAMENTAIS (com plot)
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "C03-S4 — Features Comportamentais v1 vs v2 (48 runs)")

    plot_path = PLOTS_C03 / "01_s4_recall_v1_vs_v2_w10s.png"
    if plot_path.exists():
        slide.shapes.add_picture(
            str(plot_path),
            Inches(0.3), Inches(1.4),
            Inches(7.5), Inches(5.2)
        )

    tf = add_text_box(slide, Inches(8.1), Inches(1.5), Inches(5.0), Inches(2.2),
                      "7 features comportamentais (v2):",
                      font_size=15, bold=True, color=DARK_BLUE)
    features = [
        "dst_port_entropy, dst_ip_entropy",
        "flows_per_second",
        "unanswered_ratio",
        "fwd_only_ratio",
        "small_packet_ratio",
        "syn_only_ratio",
    ]
    for f in features:
        add_paragraph(tf, f"  {f}", font_size=12, color=DARK_GRAY)

    tf2 = add_text_box(slide, Inches(8.1), Inches(4.0), Inches(5.0), Inches(2.8),
                       "Veredicto @w=10s:",
                       font_size=15, bold=True, color=DARK_BLUE)
    verdicts = [
        ("ICMP:", "v2 desbloqueia (0->50%)", GREEN),
        ("Recon:", "v2 melhor (39->45%)", GREEN),
        ("SYN:", "v1 melhor (38% vs 31%)", RED),
        ("Mirai:", "v1 melhor (46% vs 38%)", RED),
        ("TCP:", "Ambos 0%", DARK_GRAY),
        ("FPR:", "v1=2.9% -> v2=14.3%", RED),
    ]
    for label, desc, color in verdicts:
        add_paragraph(tf2, f"{label} {desc}", font_size=13, color=color)

    add_paragraph(tf2, "", font_size=6)
    add_paragraph(tf2,
        "Diagnostico: curse of dimensionality (~210 vetores, 19 features)",
        font_size=12, italic=True, color=DARK_BLUE)

    add_citation_footer(slide,
        'Zimek et al. (2012) "Unsupervised Outlier Detection in High-D Numerical Data" SADM 5(5)  |  '
        'Beyer et al. (1999) "When Is Nearest Neighbor Meaningful?" ICDT')

    # ==================================================================
    # SLIDE 12: C04 — 5 ADAPTACOES TECNICAS
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "C04 — 5 Adaptacoes Tecnicas no MicroTEDAclus")

    table = add_table(slide, Inches(0.3), Inches(1.4),
                      Inches(12.7), Inches(4.3), 6, 4)

    for i, h in enumerate(["Aspecto", "Original (Maia 2020)", "Adaptacao Propria", "Suporte Teorico"]):
        table.cell(0, i).text = h
    style_header_row(table)

    diffs = [
        ["Variancia",
         "σ² = (||δ||·2/dim)²\nSubestima ~70x em 17D",
         "Welford: dot(δ, δ')\nSoma todas dimensoes",
         "Welford (1962)\nChan et al. (1983)"],
        ["Eccentricity",
         "ξ = (||δ||·2/dim)² / (n·σ²)\nInconsistente",
         "ξ = Σ(diff²) / (n·σ²)\nConsistente com Welford",
         "Angelov (2014)"],
        ["Update policy",
         "Atualiza TODOS os\nclusters aceitantes",
         "Atualiza SO o melhor\n(maior typicality)",
         "Kohonen (1990)\nDenStream (2006)\nNS-TEDA (2024)"],
        ["Cluster n=1",
         "Sem protecao especial\n-> cluster morre rapido",
         "threshold=13 (m=5)\n-> permite crescimento",
         "Maia (2020)\ncold start"],
        ["Cluster n=2",
         "So checa var > limit",
         "Guard duplo:\nζ > thr E σ² >= r0",
         "Adaptacao propria"],
    ]
    for r, row_data in enumerate(diffs):
        for c, val in enumerate(row_data):
            bold = (r == 0)
            clr = RED if r == 0 and c == 1 else (DARK_GREEN if r == 0 and c == 2 else BLACK)
            set_cell(table, r+1, c, val, font_size=10, bold=bold, color=clr)
        if r == 0:
            table.cell(r+1, 1).fill.solid()
            table.cell(r+1, 1).fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
            table.cell(r+1, 2).fill.solid()
            table.cell(r+1, 2).fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)

    add_insight_box(slide, Inches(0.3), Inches(5.9), Inches(12.7), Inches(1.05),
        "Por que isso e contribuicao tecnica original",
        "Nenhum paper TEDA testou com >6 features. Maia (2020) validou em 2D sintetico. "
        "Aplicacao em 17D expoe limitacao estrutural previsto por Beyer (1999) e Aggarwal (2001) "
        "sobre curse of dimensionality em distancias Euclidianas.")

    add_citation_footer(slide,
        'Welford (1962) Technometrics 4(3)  |  '
        'Kohonen (1990) "Self-Organizing Map" Proc IEEE 78(9)  |  '
        'Cao et al. (2006) "DenStream" SDM  |  '
        'Chen et al. (2024) "NS-TEDA" CMC 78(2)')

    # ==================================================================
    # SLIDE 13: POR QUE A FORMULA ORIGINAL FALHA (NOVO)
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Por que a Formula Original Falha — Demonstracao Numerica")

    # A fórmula original
    shape = slide.shapes.add_shape(5, Inches(0.5), Inches(1.4),
                                   Inches(12.3), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
    shape.line.color.rgb = RED
    shape.line.width = Pt(1.5)

    add_text_box(slide, Inches(0.7), Inches(1.5), Inches(12), Inches(0.4),
                 "Formula original (Maia 2020, validada em 2D):",
                 font_size=13, bold=True, color=RED)
    add_text_box(slide, Inches(0.7), Inches(1.85), Inches(12), Inches(0.7),
                 "σ²  =  ( ||δ|| × 2/dim )²     onde δ = x − μ",
                 font_size=20, bold=True, color=BLACK,
                 font_name="Cambria Math", alignment=PP_ALIGN.CENTER)

    # Tabela dimensional
    table = add_table(slide, Inches(0.8), Inches(2.85),
                      Inches(6.5), Inches(2.8), 6, 3)
    for i, h in enumerate(["Dimensao", "Fator (2/dim)²", "Efeito"]):
        table.cell(0, i).text = h
    style_header_row(table)

    dim_data = [
        ("2D (Maia 2020)", "(2/2)² = 1.00", "Neutro — funciona", GREEN),
        ("5D", "(2/5)² = 0.16", "Subestima 6x", ORANGE),
        ("10D", "(2/10)² = 0.04", "Subestima 25x", ORANGE),
        ("17D (nosso)", "(2/17)² = 0.014", "Subestima ~70x", RED),
        ("32D", "(2/32)² = 0.004", "Subestima 250x", RED),
    ]
    for r, (dim, factor, effect, color) in enumerate(dim_data):
        set_cell(table, r+1, 0, dim, font_size=12)
        set_cell(table, r+1, 1, factor, font_size=12)
        set_cell(table, r+1, 2, effect, font_size=12, bold=True, color=color)
    # destaca linha 17D
    for c in range(3):
        table.cell(4, c).fill.solid()
        table.cell(4, c).fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)

    # Cascata de falha
    tf = add_text_box(slide, Inches(7.6), Inches(2.85), Inches(5.5), Inches(3.0),
                      "Cascata de falha em 17D:",
                      font_size=14, bold=True, color=DARK_BLUE)
    steps = [
        "1. σ² real ≈ 6.18   |   σ² original ≈ 0.087",
        "2. Razao: 71x de subestimacao",
        "3. Eccentricity ξ = ... / (n·σ²)  -> INFLADA",
        "4. Chebyshev rejeita quase todos os pontos",
        "5. Cada ponto cria novo cluster",
        "6. FPR catastrofico: 54% em trafego benigno",
    ]
    for s in steps:
        add_paragraph(tf, s, font_size=12, color=DARK_GRAY, space_before=Pt(4))

    add_paragraph(tf, "", font_size=6)
    add_paragraph(tf, "Welford corrige:",
                  font_size=14, bold=True, color=GREEN)
    add_paragraph(tf,
        "σ² incremental via produto escalar dot(δ_old, δ_new) "
        "soma contribuicoes de TODAS as dimensoes corretamente.",
        font_size=11, color=DARK_GRAY)

    add_insight_box(slide, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.95),
        "Esta e a contribuicao tecnica central da dissertacao",
        "Correcao identificada empiricamente (167 exps) e teoricamente fundamentada "
        "(Beyer 1999, Aggarwal 2001, Zimek 2012 sobre curse of dimensionality).")

    add_citation_footer(slide,
        'Beyer et al. (1999) "When Is Nearest Neighbor Meaningful?" ICDT  |  '
        'Aggarwal et al. (2001) "Surprising Behavior of Distance Metrics in High-D" ICDT  |  '
        'Chan et al. (1983) "Algorithms for Computing Sample Variance" Am Stat 37(3)')

    # ==================================================================
    # SLIDE 14: C04 — IMPACTO QUANTITATIVO (com plot)
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "C04 — Impacto Quantitativo (30 runs)")

    plot_path = PLOTS_C04 / "06_campaign04_dashboard.png"
    if plot_path.exists():
        slide.shapes.add_picture(
            str(plot_path),
            Inches(0.2), Inches(1.3),
            Inches(8.5), Inches(5.5)
        )

    tf = add_text_box(slide, Inches(8.9), Inches(1.4), Inches(4.2), Inches(2.4),
                      "FPR Benigno:",
                      font_size=18, bold=True, color=DARK_BLUE)

    fpr_data = [
        ("Flow-level", "3.9%", "54.4%", "14x"),
        ("Win v1 w=10s", "2.9%", "41.9%", "14x"),
        ("Win v1 w=30s", "5.0%", "74.5%", "15x"),
        ("Win v2 w=10s", "14.3%", "45.5%", "3x"),
    ]
    for label, ours, orig, ratio in fpr_data:
        add_paragraph(tf, f"{label}:", font_size=12, bold=True, color=DARK_GRAY)
        add_paragraph(tf,
                      f"  Propria {ours}  |  Original {orig}  ({ratio})",
                      font_size=11, color=RED)

    shape = slide.shapes.add_shape(5, Inches(8.8), Inches(4.8),
                                   Inches(4.3), Inches(2.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(1.5)

    tf2 = add_text_box(slide, Inches(9.0), Inches(4.9), Inches(3.9), Inches(2.0),
                       "Contribuicao validada:",
                       font_size=14, bold=True, color=GREEN)
    add_paragraph(tf2,
        "As 5 adaptacoes sao NECESSARIAS para MicroTEDAclus funcionar em alta dimensionalidade.",
        font_size=12, color=DARK_GRAY)
    add_paragraph(tf2,
        "Sem elas: IDS inutilizavel (FPR 42-75%). Com elas: FPR 3-15%.",
        font_size=12, bold=True, color=DARK_GREEN)

    add_citation_footer(slide,
        "Campaign-04: 30 runs, 2026-04-03. Package evolclustering (Maia 2020) vs src/detector/micro_teda.py")

    # ==================================================================
    # SLIDE 15: CONSOLIDADO — MELHOR POR ATAQUE
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Consolidado — Melhor Resultado por Ataque")

    table = add_table(slide, Inches(0.5), Inches(1.4),
                      Inches(12.3), Inches(3.3), 7, 6)

    for i, h in enumerate(["Ataque", "Campanha", "Config", "Recall", "F1", "FPR"]):
        table.cell(0, i).text = h
    style_header_row(table)

    consolidated = [
        ["DDoS-ICMP", "C03-S4", "v2 / w10s / r0=0.10", "50.0%", "5.6%", "15.7%"],
        ["DDoS-SYN", "C03-S4", "v2 / w30s / r0=0.05", "61.5%", "21.6%", "36.1%"],
        ["DDoS-TCP", "—", "Indetectavel", "0.0%", "0.0%", "—"],
        ["Mirai", "C03-S4", "v1 / w10s / r0=0.10", "46.2%", "23.1%", "15.5%"],
        ["Recon", "C03-S4", "v2 / w10s / r0=0.05", "49.1%", "43.7%", "12.9%"],
        ["Benigno (FPR)", "C02-S1", "flow-level / r0=0.10", "—", "—", "3.5%"],
    ]
    for r, row_data in enumerate(consolidated):
        for c, val in enumerate(row_data):
            bold = (r == 4)
            color = GREEN if val == "43.7%" else (RED if val in ("0.0%", "Indetectavel") else BLACK)
            set_cell(table, r+1, c, val, bold=bold, color=color, font_size=12)
    # Destaca linha Recon
    for c in range(6):
        table.cell(5, c).fill.solid()
        table.cell(5, c).fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)

    # Caixa de destaque
    shape = slide.shapes.add_shape(5, Inches(0.5), Inches(5.0),
                                   Inches(12.3), Inches(1.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(1.5)

    tf = add_text_box(slide, Inches(0.8), Inches(5.1), Inches(11.8), Inches(1.7),
                      "Destaque: Recon-PortScan F1 = 43.7%",
                      font_size=20, bold=True, color=GREEN)
    add_paragraph(tf,
        "Melhor resultado nao-supervisionado da dissertacao "
        "(Recall 49.1%, Precision 39.4%, FPR 12.9%)",
        font_size=14, color=DARK_GRAY)
    add_paragraph(tf,
        "Alternativa r0=0.15: Precision 56.7%, FPR apenas 4.2%, F1=42.0% "
        "(trade-off operacional viavel)",
        font_size=14, color=DARK_GRAY)
    add_paragraph(tf,
        "Contexto: metodos supervisionados no CICIoT2023 atingem F1>95% (requer labels). "
        "Metodos nao-supervisionados publicados no mesmo dataset sao quase inexistentes — "
        "sem baseline direto para comparacao (lacuna; trabalho futuro).",
        font_size=11, italic=True, color=DARK_BLUE)

    add_citation_footer(slide,
        'Neto et al. (2023) "CICIoT2023" Sensors 23(13)  |  '
        'Recon F1=43.7% medido em C03-S4 (v2/w10s/r0=0.05); '
        'lacuna de baseline documentada em Arp et al. (2022) para ML-security')

    # ==================================================================
    # SLIDE 16: CONTRIBUICOES CIENTIFICAS (NOVO)
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "4 Contribuicoes Cientificas")

    contributions = [
        ("1. Tecnica",
         "5 adaptacoes no MicroTEDAclus para alta dimensionalidade",
         "Correcao da formula de variancia (70x de erro em 17D), update seletivo, "
         "cold start guards. Validado em C04 com 30 runs: FPR 3.9% vs 54.4%.",
         "Primeiro estudo empirico de TEDA em dados de rede de alta dimensao.",
         MEDIUM_BLUE),
        ("2. Empirica",
         "167 experimentos documentando limites da deteccao por anomalia em IoT IDS",
         "Ablation study rigoroso confirma empiricamente as previsoes teoricas de "
         "Sommer & Paxson (2010) sobre o gap semantico anomalia-ataque.",
         "Framing como \"first controlled comparison\" no CICIoT2023.",
         GREEN),
        ("3. Metodologica",
         "Framework de \"Dimensoes Anomalas\"",
         "Modelo preditivo: detectabilidade e proporcional ao numero de dimensoes do espaco "
         "de features em que o ataque deforma o trafego. Explica hierarquia observada.",
         "Ferramenta de analise proposta a partir dos dados. Inspirada em Lakhina (2005).",
         ORANGE),
        ("4. Resultado Positivo",
         "Recon-PortScan: F1=43.7% nao-supervisionado",
         "Pipeline Kafka + MicroTEDAclus funciona para ataques com assinatura estatistica "
         "distinta (5+ dimensoes anomalas). Comparavel a baselines nao-supervisionados.",
         "Viabilidade do pipeline demonstrada no subconjunto tratavel.",
         DARK_GREEN),
    ]

    for i, (num, title, desc, framing, color) in enumerate(contributions):
        x = Inches(0.5) + Inches((i % 2) * 6.3)
        y = Inches(1.4) + Inches((i // 2) * 2.85)

        shape = slide.shapes.add_shape(5, x, y, Inches(6.1), Inches(2.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        tf = add_text_box(slide, x + Inches(0.2), y + Inches(0.1),
                          Inches(5.7), Inches(0.5),
                          num, font_size=14, bold=True, color=color)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.5),
                     Inches(5.7), Inches(0.5),
                     title, font_size=14, bold=True, color=DARK_BLUE)
        add_text_box(slide, x + Inches(0.2), y + Inches(1.05),
                     Inches(5.7), Inches(1.1),
                     desc, font_size=11, color=DARK_GRAY)
        add_text_box(slide, x + Inches(0.2), y + Inches(2.15),
                     Inches(5.7), Inches(0.5),
                     framing, font_size=10, italic=True, color=color)

    add_citation_footer(slide,
        'Matosin et al. (2014) "Negative Results" DMM 7(2)  |  '
        'Lipton & Steinhardt (2019) "Troubling Trends in ML Scholarship" Queue 17(1)')

    # ==================================================================
    # SLIDE 17: INSIGHTS / O QUE APRENDEMOS
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "O que Aprendemos — 6 Insights")

    insights = [
        ("1. Deteccao per-flow e fundamentalmente limitada",
         "Flows DDoS sao estatisticamente indistinguiveis de flows IoT benignos. "
         "Mais features per-flow (17->32) = zero impacto. Problema e estrutural.",
         "Sommer & Paxson (2010), Chen et al. (2021)"),
        ("2. Janelas temporais sao a direcao certa",
         'Mudam a pergunta do detector. Melhorias de 10-20x no Recall, mas FPR sobe.',
         "Lakhina et al. (2004, 2005), Li et al. (2023)"),
        ("3. Curse of dimensionality",
         "~210 vetores com 19 features -> clusters nao convergem. "
         "Regime de poucos dados em alta dimensionalidade e toxico.",
         "Beyer (1999), Zimek et al. (2012)"),
        ("4. Nao existe config unica otima",
         "ICMP precisa de v2, Mirai/SYN funcionam com v1. "
         "DDoS-TCP e indistinguivel em qualquer config (0 dimensoes anomalas).",
         "Zargar et al. (2013)"),
        ("5. Resultado positivo: Recon F1=43.7%",
         "Pipeline funciona para ataques com assinatura estatistica distinta. "
         "Sem baseline nao-supervisionado direto no CICIoT2023 — lacuna explicita; "
         "trabalho futuro: adaptar Isolation Forest da Fase 1 para streaming (~1-2 dias).",
         "Lacuna metodologica reconhecida"),
        ("6. Adaptacao ao dominio e contribuicao tecnica",
         "Implementacao original produz FPR 42-75% (inutilizavel). "
         "5 adaptacoes proprias reduzem para 3-15%. Contribuicao validada em C04.",
         "Welford (1962), Kohonen (1990), NS-TEDA (2024)"),
    ]
    for i, (title, desc, refs) in enumerate(insights):
        y = Inches(1.5) + Inches(i * 0.95)
        tf = add_text_box(slide, Inches(0.5), y, Inches(12.3), Inches(0.9),
                          title, font_size=15, bold=True, color=DARK_BLUE)
        add_paragraph(tf, desc, font_size=12, color=DARK_GRAY)
        add_paragraph(tf, f"Ref: {refs}", font_size=10, italic=True, color=LIGHT_GRAY)

    # ==================================================================
    # SLIDE 18: PROXIMOS PASSOS — DECISAO
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Proximos Passos — Decisao Necessaria")

    add_text_box(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
                 "Dado o prazo (~5 semanas ate defesa), qual caminho priorizar?",
                 font_size=16, bold=True, color=DARK_BLUE,
                 alignment=PP_ALIGN.CENTER)

    # Opcao A
    shape = slide.shapes.add_shape(5, Inches(0.3), Inches(2.0),
                                   Inches(4.2), Inches(4.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    shape.line.color.rgb = MEDIUM_BLUE
    shape.line.width = Pt(2)

    tf = add_text_box(slide, Inches(0.5), Inches(2.1), Inches(3.8), Inches(4.4),
                      "Opcao A", font_size=22, bold=True, color=MEDIUM_BLUE)
    add_paragraph(tf, "S5 + Escrita", font_size=16, bold=True, color=DARK_BLUE)
    add_paragraph(tf, "~1-2 semanas S5 + 3-4 escrita", font_size=12, color=DARK_GRAY)
    add_paragraph(tf, "", font_size=6)
    add_paragraph(tf, "S5 — Two-Stage Detection:", font_size=13, bold=True, color=DARK_GRAY)
    add_paragraph(tf, "Stage 1: per-flow (FPR ~3.5%)", font_size=12, color=DARK_GRAY)
    add_paragraph(tf, "Stage 2: concentracao por IP", font_size=12, color=DARK_GRAY)
    add_paragraph(tf, "", font_size=6)
    add_paragraph(tf, "+ Ataca trade-off FPR/Recall", font_size=11, color=GREEN)
    add_paragraph(tf, "+ Base teorica solida", font_size=11, color=GREEN)
    add_paragraph(tf, "− Risco de nao funcionar", font_size=11, color=RED)
    add_paragraph(tf, "− Perde 1-2 semanas se falhar", font_size=11, color=RED)

    # Opcao B
    shape = slide.shapes.add_shape(5, Inches(4.6), Inches(2.0),
                                   Inches(4.2), Inches(4.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(2)

    tf = add_text_box(slide, Inches(4.8), Inches(2.1), Inches(3.8), Inches(4.4),
                      "Opcao B", font_size=22, bold=True, color=GREEN)
    add_paragraph(tf, "Consolidar + Escrever", font_size=16, bold=True, color=DARK_BLUE)
    add_paragraph(tf, "~5 semanas inteiras de escrita", font_size=12, color=DARK_GRAY)
    add_paragraph(tf, "", font_size=6)
    add_paragraph(tf, "Foco:", font_size=13, bold=True, color=DARK_GRAY)
    add_paragraph(tf, "167 exps como contrib. metodologica", font_size=12, color=DARK_GRAY)
    add_paragraph(tf, "5 adaptacoes como contrib. tecnica", font_size=12, color=DARK_GRAY)
    add_paragraph(tf, "Negative results framing", font_size=12, color=DARK_GRAY)
    add_paragraph(tf, "", font_size=6)
    add_paragraph(tf, "+ Minimo risco de prazo", font_size=11, color=GREEN)
    add_paragraph(tf, "+ Revisao completa do orientador", font_size=11, color=GREEN)
    add_paragraph(tf, "− Menos ambicioso", font_size=11, color=ORANGE)
    add_paragraph(tf, "− Sem resultado forte adicional", font_size=11, color=ORANGE)

    # Opcao C
    shape = slide.shapes.add_shape(5, Inches(8.9), Inches(2.0),
                                   Inches(4.2), Inches(4.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
    shape.line.color.rgb = ORANGE
    shape.line.width = Pt(2)

    tf = add_text_box(slide, Inches(9.1), Inches(2.1), Inches(3.8), Inches(4.4),
                      "Opcao C", font_size=22, bold=True, color=ORANGE)
    add_paragraph(tf, "S5 + S6 + Escrita", font_size=16, bold=True, color=DARK_BLUE)
    add_paragraph(tf, "~3 semanas exps + 2 escrita", font_size=12, color=DARK_GRAY)
    add_paragraph(tf, "", font_size=6)
    add_paragraph(tf, "S5 + S6:", font_size=13, bold=True, color=DARK_GRAY)
    add_paragraph(tf, "Two-Stage + threshold adaptativo", font_size=12, color=DARK_GRAY)
    add_paragraph(tf, "", font_size=6)
    add_paragraph(tf, "+ Maximiza profundidade", font_size=11, color=GREEN)
    add_paragraph(tf, "+ Maior potencial de publicacao", font_size=11, color=GREEN)
    add_paragraph(tf, "− Margem apertada de escrita", font_size=11, color=RED)
    add_paragraph(tf, "− Alto risco de prazo", font_size=11, color=RED)

    # Minha recomendacao
    add_insight_box(slide, Inches(0.3), Inches(6.8), Inches(12.7), Inches(0.5),
        "Preferencia inicial: Opcao A",
        "S5 Two-Stage e a hipotese mais promissora (ensemble de evidencias, base teorica solida). "
        "Se nao funcionar em ~1 semana, vira Opcao B efetivamente. "
        "Quero confirmar com voces.")

    # ==================================================================
    # SLIDE 19: REFERENCIAS PRINCIPAIS (NOVO)
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Referencias Principais")

    references_left = [
        ("Fundamentacao TEDA", DARK_BLUE),
        ("Angelov (2014) \"Outside the Box\" JAMRIS 8(2)", None),
        ("Maia et al. (2020) \"Evolving Clustering Based on", None),
        ("  Mixture of Typicalities\" FGCS 106", None),
        ("Welford (1962) Technometrics 4(3)", None),
        ("Chan et al. (1983) Am Stat 37(3)", None),
        ("", None),
        ("Limites de Anomaly-based IDS", DARK_BLUE),
        ("Sommer & Paxson (2010) \"Outside the Closed", None),
        ("  World\" IEEE S&P 2010", None),
        ("Gates & Taylor (2006) NSPW", None),
        ("Axelsson (2000) \"Base-Rate Fallacy\" ACM TISSEC 3(3)", None),
        ("Chandola et al. (2009) \"Anomaly Detection: A", None),
        ("  Survey\" ACM Comp Surveys 41(3)", None),
        ("", None),
        ("DDoS e deteccao agregada", DARK_BLUE),
        ("Lakhina et al. (2004, 2005) ACM SIGCOMM", None),
        ("Zargar et al. (2013) IEEE ComSurveys 15(4)", None),
        ("Kopmann et al. (2022) \"MIDA\" ESOCC", None),
    ]

    references_right = [
        ("Curse of Dimensionality", DARK_BLUE),
        ("Beyer et al. (1999) \"When Is Nearest Neighbor", None),
        ("  Meaningful?\" ICDT", None),
        ("Aggarwal et al. (2001) \"Surprising Behavior", None),
        ("  of Distance Metrics in High-D\" ICDT", None),
        ("Zimek et al. (2012) SADM 5(5)", None),
        ("", None),
        ("Stream Clustering e Prequential", DARK_BLUE),
        ("Gama et al. (2013) \"On Evaluating Stream", None),
        ("  Learning Algorithms\" Machine Learning 90(3)", None),
        ("Kohonen (1990) \"Self-Organizing Map\" Proc IEEE", None),
        ("Cao et al. (2006) \"DenStream\" SDM", None),
        ("", None),
        ("Metodologia e Rigor", DARK_BLUE),
        ("Arp et al. (2022) \"Dos and Don'ts of ML in", None),
        ("  Computer Security\" USENIX Security", None),
        ("Pendlebury et al. (2019) \"TESSERACT\" USENIX Sec", None),
        ("", None),
        ("Dataset", DARK_BLUE),
        ("Neto et al. (2023) CICIoT2023 Sensors 23(13)", None),
    ]

    tf_left = add_text_box(slide, Inches(0.5), Inches(1.3), Inches(6.2), Inches(6.0),
                           "", font_size=11)
    for i, (text, color) in enumerate(references_left):
        if i == 0:
            tf_left.paragraphs[0].text = text
            tf_left.paragraphs[0].font.size = Pt(13)
            tf_left.paragraphs[0].font.bold = True
            tf_left.paragraphs[0].font.color.rgb = color
            tf_left.paragraphs[0].font.name = "Calibri"
        else:
            bold = color is not None
            c = color if color is not None else DARK_GRAY
            sz = 13 if bold else 11
            add_paragraph(tf_left, text, font_size=sz, bold=bold, color=c,
                          space_before=Pt(3))

    tf_right = add_text_box(slide, Inches(7.0), Inches(1.3), Inches(6.2), Inches(6.0),
                            "", font_size=11)
    for i, (text, color) in enumerate(references_right):
        if i == 0:
            tf_right.paragraphs[0].text = text
            tf_right.paragraphs[0].font.size = Pt(13)
            tf_right.paragraphs[0].font.bold = True
            tf_right.paragraphs[0].font.color.rgb = color
            tf_right.paragraphs[0].font.name = "Calibri"
        else:
            bold = color is not None
            c = color if color is not None else DARK_GRAY
            sz = 13 if bold else 11
            add_paragraph(tf_right, text, font_size=sz, bold=bold, color=c,
                          space_before=Pt(3))

    # Save
    prs.save(str(OUTPUT_PATH))
    print(f"Apresentacao salva em: {OUTPUT_PATH}")
    print(f"Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
