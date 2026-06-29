#!/usr/bin/env python3
"""Generate the .pptx for the paper-defense meeting (2026-05-08).

10-slide deck mirroring `presentation.md`. Markdown is the source of truth;
this file produces a projectable .pptx synced to it.

Usage:
    python docs/meeting/2026-05-08-paper-defense/generate_pptx.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Paths
SCRIPT_DIR = Path(__file__).parent
OUTPUT_PATH = SCRIPT_DIR / "2026-05-08-paper-defense.pptx"

# Colors (consistent with prior meetings)
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MEDIUM_BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
LIGHT_GRAY = RGBColor(0x90, 0x90, 0x90)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
GOLD = RGBColor(0xB8, 0x86, 0x0B)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, *,
                 font_size=18, bold=False, italic=False, color=BLACK,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, *, font_size=18, bold=False, italic=False,
                  color=BLACK, alignment=PP_ALIGN.LEFT,
                  space_before=Pt(6), font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_bullets(slide, left, top, width, height, items, *,
                font_size=16, color=BLACK, bullet="• "):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet + item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_before = Pt(4)
    return tf


def add_title(slide, text, *, color=DARK_BLUE):
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.8),
                 text, font_size=28, bold=True, color=color)


def add_footer(slide, text, color=LIGHT_GRAY):
    add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12.5), Inches(0.4),
                 text, font_size=10, color=color, alignment=PP_ALIGN.RIGHT)


# ── Slide builders ─────────────────────────────────────────────

def build_slide_1_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, DARK_BLUE)

    add_text_box(s, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.2),
                 "Phase Transition Characterization in MicroTEDAclus",
                 font_size=40, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(s, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.6),
                 "Defesa do paper SoftCom 2026",
                 font_size=22, italic=True, color=LIGHT_BLUE,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.5),
                 "Augusto Custódio Vicente",
                 font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.5),
                 "Orientador: Frederico Gadelha Guimarães",
                 font_size=16, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

    add_text_box(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
                 "PPGEE — UFMG · 2026-05-08 · Submission deadline 2026-05-11",
                 font_size=14, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)


def build_slide_2_pivot(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, WHITE)
    add_title(s, "1. Pivô do Framing desde 17/04")

    add_text_box(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.6),
                 "Reunião anterior: framing operacional. Hoje: framing estrutural.",
                 font_size=16, italic=True, color=DARK_GRAY)

    add_bullets(s, Inches(0.7), Inches(2.1), Inches(12), Inches(4.5), [
        "Framing 17/04: \"5 adaptações para corrigir MicroTEDAclus em IoT IDS\". "
        "Defensável mas modesto.",
        "Análise da Campaign-05 (22-28/04): 3 regimes patológicos com topologias estruturalmente distintas.",
        "Investigação 28/04 → 04/05: identificado o mecanismo causal — comparação σ² vs r₀ no denominador da eccentricidade.",
        "Experimento controlado (Exp 3, 04/05, 1.620 runs): predição estrutural λ* ∝ √r₀ confirmada com check_ratio = 1,00 EXATO.",
        "Pivô: \"Phase transition characterization\" como contribuição teórica central.",
    ], font_size=18)

    add_text_box(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
                 "→ Pergunta central desta reunião: o framing regime change é "
                 "defensável vs orientador e reviewers SoftCom?",
                 font_size=16, bold=True, color=MEDIUM_BLUE)

    add_footer(s, "Slide 1 / 11")


def build_slide_3_mechanism(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, WHITE)
    add_title(s, "2. O Mecanismo Estrutural")

    add_text_box(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
                 "Em corrected.py:233 — uma linha define dois regimes operacionais.",
                 font_size=16, italic=True, color=DARK_GRAY)

    add_text_box(s, Inches(1.0), Inches(2.0), Inches(11), Inches(0.7),
                 "effective_var = max(var, self.r0)",
                 font_size=24, bold=True, color=GOLD,
                 font_name="Consolas")

    add_bullets(s, Inches(0.7), Inches(3.0), Inches(12), Inches(3.5), [
        "Função max(·) é não-suave → trava em um dos dois argumentos. Define DOIS regimes:",
        "    r0-bounded (σ² ≪ r₀): denominador = r₀ constante. Limiar de aceitação NÃO-ADAPTATIVO.",
        "    data-bounded (σ² ≫ r₀): denominador = σ². Limiar ADAPTATIVO à dispersão local.",
        "Não é gradiente; é phase transition operacional. Fronteira nítida em σ² = r₀.",
        "4 code paths em corrected.py (linhas 233, 280, 292, 296) executam essa comparação.",
        "V0 (literal Maia) tem comparação análoga em _is_outlier:200 (var > r₀) — mesmo regime structure.",
    ], font_size=16)

    add_text_box(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.5),
                 "→ Regime change é estrutural ao MicroTEDAclus, "
                 "não peculiaridade da nossa correção.",
                 font_size=15, bold=True, color=GREEN)

    add_footer(s, "Slide 2 / 11")


def build_slide_4_prediction(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, WHITE)
    add_title(s, "3. Predição Algébrica Falsificável")

    add_text_box(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
                 "Para Gaussiana 𝒩(0, λ²I_d), Welford → trace covariance → previsão de λ*",
                 font_size=15, italic=True, color=DARK_GRAY)

    add_text_box(s, Inches(0.7), Inches(2.1), Inches(12), Inches(0.6),
                 "σ² → d · λ²    (variance Welford = trace cov, confirmado em test_welford_variance.py)",
                 font_size=16, color=DARK_GRAY)

    add_text_box(s, Inches(2.5), Inches(3.0), Inches(8), Inches(0.9),
                 "λ* = √(r₀ / d)",
                 font_size=36, bold=True, color=DARK_BLUE,
                 alignment=PP_ALIGN.CENTER, font_name="Cambria Math")

    # Predicted values table
    add_text_box(s, Inches(2), Inches(4.3), Inches(9), Inches(0.5),
                 "Predições para d = 17 (CICIoT2023 v1):",
                 font_size=16, bold=True, color=BLACK)
    add_text_box(s, Inches(2.5), Inches(4.9), Inches(8), Inches(2.0),
                 "  r₀ = 10⁻³  →  λ* = 0,0077\n"
                 "  r₀ = 10⁻¹  →  λ* = 0,0767\n"
                 "  r₀ = 10⁰   →  λ* = 0,243",
                 font_size=18, font_name="Consolas", color=BLACK)

    add_text_box(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.5),
                 "→ Falsificável: se a estrutura √r₀ não emerge empiricamente, "
                 "hipótese está errada.",
                 font_size=15, bold=True, color=ORANGE)

    add_footer(s, "Slide 3 / 11")


def build_slide_4_2overd(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, WHITE)
    add_title(s, "4. A Ponte Teoria→Prática: o Fator (2/d)²")

    add_text_box(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5),
                 "A predição λ* = √(r₀/d) vale para σ² 'real'. Que σ² cada implementação acumula?",
                 font_size=14, italic=True, color=DARK_GRAY)

    # Two algorithms compared
    add_text_box(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
                 "V7 (Welford raw):  σ² → d · λ²  (correto)",
                 font_size=18, bold=True, color=GREEN, font_name="Consolas")
    add_text_box(s, Inches(0.7), Inches(2.6), Inches(12), Inches(0.5),
                 "V0 (Maia literal): σ² = (‖δ‖·2/d)² / (n-1) ≈ (4/d)·λ²  ← fator (2/d)² encolhe σ² em ~d²/4",
                 font_size=18, bold=True, color=RED, font_name="Consolas")

    # Table of effect by dimension
    table_data = [
        ("d", "(2/d)²", "Encolhimento", "Resultado V0"),
        ("2 (Maia validou)", "1", "nenhum", "OK — funciona como V7"),
        ("17 (IoT real)", "0,014", "72×", "σ² próx. de r₀ → fragmenta"),
        ("50", "0,0016", "625×", "σ² ≪ r₀ → r0-bounded"),
    ]
    table = s.shapes.add_table(4, 4, Inches(1.5), Inches(3.5),
                               Inches(10.3), Inches(2.0)).table
    for i, row in enumerate(table_data):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.bold = (i == 0)
                p.alignment = PP_ALIGN.CENTER

    add_text_box(s, Inches(0.7), Inches(5.8), Inches(12), Inches(0.6),
                 "→ (2/d)² NÃO é o \"bug central\". É o mecanismo específico que "
                 "joga V0 na fronteira de regime em alta-d.",
                 font_size=15, bold=True, color=DARK_BLUE)

    add_text_box(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
                 "Em d≤6 V0 funciona em data-bounded; em d≥10 V0 atravessa para r0-bounded paranoia.",
                 font_size=14, italic=True, color=DARK_GRAY)

    add_footer(s, "Slide 4 / 11")


def build_slide_5_hypotheses(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, LIGHT_BLUE)
    add_title(s, "5. Hipóteses Pré-Registradas")

    add_text_box(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
                 "Registradas no plano em 04/05 ANTES do experimento rodar.",
                 font_size=16, italic=True, color=DARK_GRAY)

    items = [
        ("H1", "Estrutura V7", "V7 transiciona em λ* = √(r₀/d), tolerância 2× (testa coeficiente)."),
        ("H2", "V0 difere", "V0 (sem max guard) difere qualitativamente: Cohen's d > 0,8 vs V7."),
        ("H3", "Escala universal", "Razão λ*(r₀_a)/λ*(r₀_b) preserva √(r₀_a/r₀_b) (testa estrutura)."),
    ]

    y = Inches(2.3)
    for tag, title, desc in items:
        add_text_box(s, Inches(0.7), y, Inches(1.2), Inches(0.6),
                     tag, font_size=24, bold=True, color=DARK_BLUE)
        add_text_box(s, Inches(1.9), y, Inches(2.5), Inches(0.6),
                     title, font_size=18, bold=True, color=BLACK)
        add_text_box(s, Inches(4.5), y, Inches(8.3), Inches(0.6),
                     desc, font_size=15, color=DARK_GRAY)
        y += Inches(1.0)

    add_text_box(s, Inches(0.7), Inches(5.7), Inches(12), Inches(0.6),
                 "Critério de robustez (todos simultâneos): Friedman p<0,001, "
                 "λ* dentro de 2×, razão preservada, Cohen's d > 0,8.",
                 font_size=14, color=DARK_GRAY)

    add_text_box(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.6),
                 "→ Caso B do plano: estrutura confirma com folga, "
                 "coeficiente refina (explicação em §VI Discussion).",
                 font_size=15, bold=True, color=GREEN)

    add_footer(s, "Slide 5 / 11")


def build_slide_6_setup(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, WHITE)
    add_title(s, "6. Setup Experimental (Exp 3)")

    add_bullets(s, Inches(0.7), Inches(1.5), Inches(12), Inches(5), [
        "1.620 runs: λ × r₀ × algoritmo × seed",
        "    9 valores λ log-espaçados em [10⁻³, 10¹] — cobre 4 ordens de grandeza",
        "    3 valores r₀ ∈ {10⁻³, 10⁻¹, 10⁰} — cobre 3 ordens",
        "    2 algoritmos: V0 (literal Maia) e V7 (corrigido com 5 adaptações)",
        "    30 seeds por condição (padrão do projeto: Exp 1 e Exp 2)",
        "Dimensão d = 17 (matches IoT v1 features), 1.000 amostras/run, anomalias 5σ",
        "Estatística: Friedman + Nemenyi, ANOVA + Tukey HSD em paralelo, bootstrap CI 95%",
        "Tempo de execução: ~30 min em Mac M2",
        "Código: experiments/teda-high-dim/experiments/exp03_regime_transition.py",
    ], font_size=16)

    add_footer(s, "Slide 6 / 11")


def build_slide_7_results(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, WHITE)
    add_title(s, "7. Resultado Principal — H3 confirmada exatamente")

    add_text_box(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
                 "Razões empíricas batem √(r₀_a/r₀_b) com check_ratio = 1,00 EXATO em 3/3.",
                 font_size=16, italic=True, color=DARK_GRAY)

    # Table of ratios
    add_text_box(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
                 "Razões λ* observadas vs. preditas (V7):",
                 font_size=15, bold=True, color=BLACK)

    table_data = [
        ("Comparação", "Observado", "Predito √(r₀_a/r₀_b)", "Check ratio"),
        ("λ*(0,1) / λ*(10⁻³)", "10,000", "10,000", "1,00"),
        ("λ*(1) / λ*(10⁻³)", "31,623", "31,623", "1,00"),
        ("λ*(1) / λ*(0,1)", "3,162", "3,162", "1,00"),
    ]
    table = s.shapes.add_table(4, 4, Inches(1.5), Inches(2.6),
                               Inches(10), Inches(1.7)).table
    for i, row in enumerate(table_data):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.bold = (i == 0)
                p.alignment = PP_ALIGN.CENTER

    add_bullets(s, Inches(0.7), Inches(4.6), Inches(12), Inches(2.3), [
        "Coeficiente: V7 empírico = 0,092·√r₀ vs predito 0,243·√r₀. Fator 0,38 CONSTANTE em 3 ordens de r₀.",
        "Friedman χ² ≥ 160, p ≤ 10⁻³⁰ em todas as condições — separação estatística massiva.",
        "Offset 0,38 explicado: regime indicator usa σ² por cluster; quando V7 fragmenta, σ² < trace covariance total.",
        "Anomalias 5σ inflam σ² pré-fragmentação. Refinamento quantitativo do prefactor; estrutura √r₀ é EXATA.",
    ], font_size=14)

    add_footer(s, "Slide 7 / 11")


def build_slide_8_v0_vs_v7(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, WHITE)
    add_title(s, "8. V0 vs V7 — Diferença Qualitativa Massiva")

    add_text_box(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
                 "Cohen's d até 1.376 — três ordens de grandeza acima do threshold de \"large effect\".",
                 font_size=15, italic=True, color=DARK_GRAY)

    table_data = [
        ("r₀", "λ", "FPR V0", "FPR V7", "Cohen's d"),
        ("10⁻³", "0,316", "0,998", "0,001", "+1.247"),
        ("10⁻³", "1,0",   "0,998", "0,001", "+1.376"),
        ("10⁻¹", "3,16",  "0,998", "0,001", "+1.247"),
        ("10⁰",  "10",    "0,998", "0,001", "+1.247"),
    ]
    table = s.shapes.add_table(5, 5, Inches(2), Inches(2.0),
                               Inches(9), Inches(2.5)).table
    for i, row in enumerate(table_data):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.bold = (i == 0)
                p.alignment = PP_ALIGN.CENTER

    add_bullets(s, Inches(0.7), Inches(4.7), Inches(12), Inches(2.0), [
        "Threshold Cohen 1988 \"large effect\": |d| > 0,8. Observamos magnitudes 1.000× maior.",
        "V0 colapsa em 1.000 clusters (cap) em λ ≥ 1; V7 estável em ~6 clusters em todo λ.",
        "Visual em fig_regime_v0_vs_v7.pdf: trade-off em FPR, # clusters, top-1 fraction.",
    ], font_size=15)

    add_text_box(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
                 "→ Mensagem: as 5 adaptações não \"calibram\" o detector. "
                 "Mudam-no estruturalmente.",
                 font_size=15, bold=True, color=GREEN)

    add_footer(s, "Slide 8 / 11")


def build_slide_9_maia(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, WHITE)
    add_title(s, "9. Reinterpretação de Maia 2020 — Escopo Retrospectivo")

    add_text_box(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
                 "Não refutamos Maia. Delimitamos escopo de aplicabilidade.",
                 font_size=16, italic=True, color=DARK_GRAY)

    add_bullets(s, Inches(0.7), Inches(2.0), Inches(12), Inches(4.5), [
        "Maia 2020 reporta \"r₀ = 0,001 robusto para todos os datasets\".",
        "Datasets do paper: d ∈ {2, 3}, features escala λ ~ 1 (ST-D1, ST-D2, Cassini, RBF).",
        "Substituindo na nossa fórmula: λ* = √(0,001/2) ≈ 0,022 ≪ 1. → Maia opera firmemente em data-bounded.",
        "A \"robustez de r₀\" de Maia é estabilidade DENTRO de um regime — não generalização ENTRE regimes.",
        "Em IoT com features raw e d = 17: a calibração default cai na fronteira → \"robustez\" colapsa.",
        "Politicamente: Frederico é co-autor; framing aqui é EXTENSÃO crítica de aplicabilidade, não refutação.",
    ], font_size=15)

    add_text_box(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
                 "→ Decisão: enviar email para Maia (via Frederico) com framing acima ANTES da submissão.",
                 font_size=15, bold=True, color=ORANGE)

    add_footer(s, "Slide 9 / 11")


def build_slide_10_iot_table(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, WHITE)
    add_title(s, "10. IoT Manifestation + Tabela VIII Honesta")

    add_text_box(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5),
                 "C04 per-flow: V0 FPR = 54,4% → V7 FPR = 3,9% (14× improvement). C05 baselines abaixo:",
                 font_size=14, italic=True, color=DARK_GRAY)

    table_data = [
        ("Algoritmo", "Tipo", "FPR (%)", "F1 (%)", "Tput (fl/s)", "Cal. gap (pp)"),
        ("HST", "Streaming", "47,0", "50,4", "134,6", "7,8"),
        ("V7 (ours)", "Incremental", "3,8", "9,8", "124,9", "40,6"),
        ("V0", "Incremental", "48,7", "44,8", "16,4", "0,6"),
        ("LOF", "Streaming", "19,2", "14,1", "1,7", "38,5"),
    ]
    table = s.shapes.add_table(5, 6, Inches(0.7), Inches(2.0),
                               Inches(11.9), Inches(2.6)).table
    for i, row in enumerate(table_data):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.bold = (i == 0 or (i == 2))  # bold V7 row
                p.alignment = PP_ALIGN.CENTER

    add_bullets(s, Inches(0.7), Inches(4.8), Inches(12), Inches(2.0), [
        "HST domina F1. Não escondemos. Calibration gap V7 = 40,6 pp expõe silent collapse.",
        "Reframe: V7 é caracterização teórica, não competidor SOTA. HST é referência streaming, não competidor.",
        "Campaign-06 em preparação (NÃO rodada): testa se features normalizadas colapsam V0/V7 em IoT real.",
    ], font_size=14)

    add_footer(s, "Slide 10 / 11")


def build_slide_11_close(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, DARK_BLUE)
    add_title(s, "11. Conclusão e Pedido", color=WHITE)

    add_text_box(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
                 "Paper pronto: 587 linhas, 22 referências (4 novas), 6 figuras (3 novas), 0 placeholders.",
                 font_size=16, italic=True, color=LIGHT_BLUE)

    add_text_box(s, Inches(0.7), Inches(2.2), Inches(12), Inches(0.5),
                 "Pendências antes da submissão (10/05):",
                 font_size=18, bold=True, color=WHITE)

    add_bullets(s, Inches(0.7), Inches(2.8), Inches(12), Inches(2.0), [
        "Build LaTeX limpo no Overleaf (hoje/amanhã)",
        "Verificação final dos números da Tabela VIII contra metrics_summary.csv",
        "Email para Maia, via Frederico (framing escopo retrospectivo)",
        "Sua revisão deste paper",
    ], font_size=16, color=WHITE)

    add_text_box(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.5),
                 "Perguntas diretas:",
                 font_size=18, bold=True, color=WHITE)

    add_bullets(s, Inches(0.7), Inches(5.5), Inches(12), Inches(2.0), [
        "O framing regime change é defensável vs orientador e reviewers SoftCom?",
        "Tabela VIII honesta com HST F1 > V7 está OK?",
        "OK email para Maia? Posso redigir, você revisa.",
    ], font_size=16, color=LIGHT_BLUE)

    add_footer(s, "Slide 11 / 11 · Submissão: 10/05 · Hard deadline: 11/05",
               color=LIGHT_BLUE)


# ── Main ───────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    builders = [
        build_slide_1_title,
        build_slide_2_pivot,
        build_slide_3_mechanism,
        build_slide_4_prediction,
        build_slide_4_2overd,
        build_slide_5_hypotheses,
        build_slide_6_setup,
        build_slide_7_results,
        build_slide_8_v0_vs_v7,
        build_slide_9_maia,
        build_slide_10_iot_table,
        build_slide_11_close,
    ]
    for b in builders:
        b(prs)

    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
