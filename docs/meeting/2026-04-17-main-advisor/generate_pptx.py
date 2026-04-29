#!/usr/bin/env python3
"""
Gera apresentacao PowerPoint para reuniao com orientador principal (17/abr/2026).
Uso: python docs/meeting/2026-04-17-main-advisor/generate_pptx.py

Estrutura (9 slides):
 1. Contexto e Objetivo (dark blue bg)
 2. Pipeline Streaming: Por Que Alta Dimensao Importa (white bg)
 3. TEDA: Framework Base (light blue bg)
 4. MicroTEDAclus (Maia 2020) (white bg)
 5. O Bug: Auto-Cancelamento + 3 Falhas (white bg, KEY SLIDE)
 6. As 5 Adaptacoes (white bg)
 7. Evidencia: C04 + Resultados + Gap Semantico (white bg)
 8. Experimentos a Executar (white bg)
 9. Decisoes para o Orientador (dark blue bg)
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Paths
SCRIPT_DIR = Path(__file__).parent
PROJECT_ROOT = SCRIPT_DIR.parent.parent.parent
OUTPUT_PATH = SCRIPT_DIR / "2026-04-17-main-advisor.pptx"

# Colors (same palette as the 2026-03-19 meeting script)
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MEDIUM_BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
LIGHT_GRAY = RGBColor(0x90, 0x90, 0x90)
GREEN = RGBColor(0x27, 0xAE, 0x60)
DARK_GREEN = RGBColor(0x1B, 0x5E, 0x20)
RED = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
GOLD = RGBColor(0xB8, 0x86, 0x0B)

# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ============================================================
# HELPER FUNCTIONS
# ============================================================

def set_slide_bg(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", italic=False):
    """Add a text box to slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=18, bold=False, color=BLACK,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6), italic=False):
    """Add paragraph to existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_table(slide, left, top, width, height, rows, cols):
    """Add table to slide."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table


def style_header_row(table, color=DARK_BLUE, font_color=WHITE):
    """Style the header row of a table."""
    for i, cell in enumerate(table.rows[0].cells):
        cell.fill.solid()
        cell.fill.fore_color.rgb = color
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = font_color
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.name = "Calibri"


def style_data_cell(cell, font_size=13, bold=False, color=BLACK,
                    alignment=PP_ALIGN.CENTER, italic=False):
    """Style a data cell."""
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def set_cell(table, row, col, text, font_size=13, bold=False, color=BLACK,
             italic=False, alignment=PP_ALIGN.CENTER):
    """Set cell text with styling."""
    cell = table.cell(row, col)
    cell.text = text
    style_data_cell(cell, font_size=font_size, bold=bold, color=color,
                    italic=italic, alignment=alignment)


def add_title_bar(slide, title_text):
    """Add a colored title bar at the top."""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.9),
                 title_text, font_size=30, bold=True, color=WHITE)


def add_citation_footer(slide, citation_text):
    """Add small citation text at the bottom of a slide."""
    add_text_box(slide, Inches(0.5), Inches(7.05), Inches(12.5), Inches(0.4),
                 citation_text, font_size=10, color=LIGHT_GRAY, italic=True)


def add_insight_box(slide, left, top, width, height, title, body, color=MEDIUM_BLUE):
    """Add a highlighted insight callout box."""
    shape = slide.shapes.add_shape(5, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)

    tf = add_text_box(slide, left + Inches(0.2), top + Inches(0.1),
                      width - Inches(0.4), height - Inches(0.2),
                      title, font_size=14, bold=True, color=color)
    add_paragraph(tf, body, font_size=12, color=DARK_GRAY)


# ============================================================
# PRESENTATION BUILDER
# ============================================================

def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]

    # ==================================================================
    # SLIDE 1: CONTEXTO E OBJETIVO (dark blue bg)
    # ==================================================================
    print("Gerando slide 1: Contexto e Objetivo...")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BLUE)

    add_text_box(slide, Inches(1), Inches(1.3), Inches(11), Inches(1.5),
                 "MicroTEDAclus em Alta Dimensao",
                 font_size=42, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(2.7), Inches(11), Inches(1.0),
                 "Progresso e Proximos Passos",
                 font_size=26, color=LIGHT_BLUE,
                 alignment=PP_ALIGN.CENTER)

    # Badge
    shape = slide.shapes.add_shape(5, Inches(3.0), Inches(4.2),
                                   Inches(7.3), Inches(1.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MEDIUM_BLUE
    shape.line.fill.background()

    tf = add_text_box(slide, Inches(3.2), Inches(4.3), Inches(6.9), Inches(0.6),
                      "Fase 2: 167 experimentos | 4 campanhas | defesa ~agosto 2026",
                      font_size=18, bold=True, color=WHITE,
                      alignment=PP_ALIGN.CENTER)
    add_paragraph(tf,
                  "Foco: contribuicao tecnica + decisao dos proximos experimentos",
                  font_size=16, color=RGBColor(0xBB, 0xDE, 0xFB),
                  alignment=PP_ALIGN.CENTER, space_before=Pt(10))

    tf = add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.8),
                      "Augusto | Mestrado PPGEE - UFMG | Orientacao 17/abr/2026",
                      font_size=16, color=LIGHT_BLUE,
                      alignment=PP_ALIGN.CENTER)

    # ==================================================================
    # SLIDE 2: TEDA FRAMEWORK BASE (light blue bg)
    # ==================================================================
    print("Gerando slide 2: TEDA Framework Base...")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, LIGHT_BLUE)
    add_title_bar(slide, "TEDA \u2014 Tipicality and Eccentricity Data Analytics (Angelov 2014)")

    # LEFT: Equations + definitions
    left_w = Inches(6.3)

    # Eccentricity box
    shape = slide.shapes.add_shape(5, Inches(0.3), Inches(1.5),
                                   left_w, Inches(1.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = DARK_BLUE
    shape.line.width = Pt(2)

    add_text_box(slide, Inches(0.5), Inches(1.55), left_w - Inches(0.4), Inches(0.4),
                 "Eccentricidade \u2014 qu\u00e3o diferente \u00e9 o ponto x:",
                 font_size=13, bold=True, color=DARK_BLUE)
    add_text_box(slide, Inches(0.5), Inches(1.95), left_w - Inches(0.4), Inches(0.5),
                 "\u03be(x) = 1/k + \u2016x \u2212 \u03bc\u2016\u00b2 / (k \xb7 \u03c3\u00b2)",
                 font_size=22, bold=True, color=DARK_BLUE,
                 font_name="Cambria Math", alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(2.5), left_w - Inches(0.4), Inches(0.4),
                 "\u03be \u2248 1/k \u2192 t\u00edpico (perto da m\u00e9dia)  |  \u03be \u226b 1/k \u2192 outlier (longe)",
                 font_size=11, italic=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # Typicality box
    shape = slide.shapes.add_shape(5, Inches(0.3), Inches(3.3),
                                   left_w, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(2)

    add_text_box(slide, Inches(0.5), Inches(3.35), left_w - Inches(0.4), Inches(0.4),
                 "Tipicalidade \u2014 complemento (qu\u00e3o t\u00edpico):",
                 font_size=13, bold=True, color=DARK_GREEN)
    add_text_box(slide, Inches(0.5), Inches(3.75), left_w - Inches(0.4), Inches(0.5),
                 "\u03c4(x) = 1 \u2212 \u03be(x)",
                 font_size=22, bold=True, color=DARK_GREEN,
                 font_name="Cambria Math", alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(4.2), left_w - Inches(0.4), Inches(0.3),
                 "\u03c4 \u2248 1 \u2192 muito t\u00edpico  |  \u03c4 \u2264 0 \u2192 outlier",
                 font_size=11, italic=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # RIGHT: Pseudocode TEDA
    right_x = Inches(6.8)
    right_w = Inches(6.3)

    shape = slide.shapes.add_shape(5, right_x, Inches(1.5),
                                   right_w, Inches(3.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    shape.line.color.rgb = DARK_GRAY
    shape.line.width = Pt(1)

    add_text_box(slide, right_x + Inches(0.15), Inches(1.55),
                 right_w - Inches(0.3), Inches(0.4),
                 "Pseudoc\u00f3digo TEDA (centro \u00fanico):",
                 font_size=13, bold=True, color=DARK_BLUE)

    pseudo_teda = (
        "para cada ponto x:\n"
        "  \u03bc \u2190 \u03bc + (x \u2212 \u03bc) / k\n"
        "  \u03c3\u00b2 \u2190 atualizar via Welford\n"
        "  \u03be \u2190 1/k + \u2016x\u2212\u03bc\u2016\u00b2 / (k\xb7\u03c3\u00b2)\n"
        "  \u03b6 \u2190 \u03be / 2    (normalizada)\n"
        "  se \u03b6 > (m\u00b2+1)/(2k):   \u2192 OUTLIER\n"
        "  sen\u00e3o:                \u2192 normal, atualizar"
    )
    add_text_box(slide, right_x + Inches(0.15), Inches(2.0),
                 right_w - Inches(0.3), Inches(2.5),
                 pseudo_teda, font_size=12, color=DARK_GRAY,
                 font_name="Consolas")

    # Properties bar (compact, 3 cols)
    props = [
        ("N\u00e3o-param\u00e9trico", MEDIUM_BLUE),
        ("O(1) por ponto (Welford)", GREEN),
        ("Chebyshev: P(|X\u2212\u03bc|\u2265m\u03c3)\u22641/m\u00b2", ORANGE),
    ]
    for i, (txt, col) in enumerate(props):
        px = Inches(0.3) + Inches(i * 4.35)
        shape = slide.shapes.add_shape(5, px, Inches(4.7),
                                       Inches(4.2), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = col
        shape.line.width = Pt(1.5)
        add_text_box(slide, px + Inches(0.15), Inches(4.78),
                     Inches(3.9), Inches(0.4),
                     txt, font_size=12, bold=True, color=col,
                     alignment=PP_ALIGN.CENTER)

    # Footer limitation
    shape = slide.shapes.add_shape(5, Inches(0.3), Inches(5.6),
                                   Inches(12.8), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
    shape.line.color.rgb = RED
    shape.line.width = Pt(1.5)

    add_text_box(slide, Inches(0.6), Inches(5.7), Inches(12.2), Inches(0.9),
                 "Limita\u00e7\u00e3o: centro \u00fanico \u2192 outliers contaminam \u03bc e \u03c3\u00b2 \u2192 motiva MicroTEDAclus",
                 font_size=15, bold=True, color=RED)

    # ==================================================================
    # SLIDE 3: MICROTEDACLUS (white bg)
    # ==================================================================
    print("Gerando slide 3: MicroTEDAclus...")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "MicroTEDAclus \u2014 M\u00faltiplos Micro-Clusters (Maia et al. 2020)")

    # LEFT: Key concepts
    left_w = Inches(6.0)

    concepts = [
        ("Cada micro-cluster mant\u00e9m (\u03bc\u1d62, \u03c3\u00b2\u1d62, S\u1d62) isoladamente",
         "Estat\u00edsticas independentes \u2192 outliers n\u00e3o contaminam outros clusters"),
        ("K autom\u00e1tico, concept drift nativo",
         "Novos padr\u00f5es criam novos clusters; K n\u00e3o \u00e9 par\u00e2metro"),
        ("m(k) din\u00e2mico: protege clusters jovens",
         "k=1: m\u22480.6 (permissivo) | k\u2192\u221e: m\u21923 (estrito, 89% Chebyshev)"),
        ("\u00danico par\u00e2metro: r\u2080 (piso de vari\u00e2ncia)",
         "Maia 2020: r\u2080=0.001 fixo. N\u00f3s descobrimos: depende da escala dos dados"),
    ]
    for i, (main, sub) in enumerate(concepts):
        y = Inches(1.6) + Inches(i * 1.1)
        tf = add_text_box(slide, Inches(0.5), y, left_w, Inches(0.4),
                          main, font_size=15, bold=True, color=DARK_BLUE)
        add_text_box(slide, Inches(0.7), y + Inches(0.35), left_w - Inches(0.2), Inches(0.4),
                     sub, font_size=12, color=DARK_GRAY, italic=True)

    # RIGHT: Pseudocode MicroTEDAclus
    right_x = Inches(6.5)
    right_w = Inches(6.5)

    shape = slide.shapes.add_shape(5, right_x, Inches(1.5),
                                   right_w, Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    shape.line.color.rgb = DARK_GRAY
    shape.line.width = Pt(1)

    add_text_box(slide, right_x + Inches(0.15), Inches(1.55),
                 right_w - Inches(0.3), Inches(0.4),
                 "Pseudoc\u00f3digo MicroTEDAclus:",
                 font_size=13, bold=True, color=DARK_BLUE)

    pseudo_micro = (
        "para cada ponto x:\n"
        "  se primeiro ponto:\n"
        "    criar MC\u2081(\u03bc=x, \u03c3\u00b2=0, S=1)\n"
        "\n"
        "  para cada MC\u1d62:\n"
        "    S\u1d62 \u2190 S\u1d62 + 1\n"
        "    \u03bc \u2190 \u03bc + (x\u2212\u03bc)/S       # m\u00e9dia incremental\n"
        "    \u03c3\u00b2 \u2190 atualizar(\u03b4)    # \u2190 BUG AQUI\n"
        "    \u03be \u2190 1/S + \u2016x\u2212\u03bc\u2016\u00b2/(S\u00b7\u03c3\u00b2)\n"
        "    \u03b6 \u2190 \u03be/2\n"
        "\n"
        "    se S < 3:              # \u2190 PATH 2\n"
        "      outlier \u2190 (\u03c3\u00b2 > r\u2080)\n"
        "    sen\u00e3o:               # Chebyshev OK\n"
        "      outlier \u2190 (\u03b6 > thr)\n"
        "\n"
        "  se algum aceita:\n"
        "    best \u2190 MC com max \u03c4\u1d62(x)\n"
        "    atualizar best: \u03bc, \u03c3\u00b2, S\n"
        "    life += (\u221a\u03c3\u00b2\u2212dist)/\u221a\u03c3\u00b2  # \u2190 PATH 3\n"
        "    merge se dist < 2(\u221a\u03c3\u00b2) # \u2190 PATH 1\n"
        "    \u2192 NORMAL\n"
        "\n"
        "  sen\u00e3o (rejeitado por todos):\n"
        "    criar MC\u2099\u208a\u2081(\u03bc=x, \u03c3\u00b2=0, S=1)\n"
        "    \u2192 ANOMALIA"
    )
    add_text_box(slide, right_x + Inches(0.15), Inches(2.0),
                 right_w - Inches(0.3), Inches(3.8),
                 pseudo_micro, font_size=11, color=DARK_GRAY,
                 font_name="Consolas")

    # Footer in red
    shape = slide.shapes.add_shape(5, Inches(0.3), Inches(6.2),
                                   Inches(12.8), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
    shape.line.color.rgb = RED
    shape.line.width = Pt(1.5)

    add_text_box(slide, Inches(0.6), Inches(6.3), Inches(12.2), Inches(0.7),
                 "Validado em 2\u20136D. Nunca testado em alta dimens\u00e3o (nosso caso: d=17\u201319).",
                 font_size=15, bold=True, color=RED)

    # ==================================================================
    # SLIDE 4: O BUG (white bg, KEY SLIDE)
    # ==================================================================
    print("Gerando slide 4: O Bug Dimensional (slide chave)...")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "O Bug Dimensional: Onde Funciona (Por Acidente) e Onde Colapsa")

    # LEFT half: green-tinted box
    left_x = Inches(0.3)
    box_top = Inches(1.5)
    half_width = Inches(6.3)
    box_h = Inches(4.5)

    shape = slide.shapes.add_shape(5, left_x, box_top, half_width, box_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(2)

    tf = add_text_box(slide, left_x + Inches(0.2), box_top + Inches(0.1),
                      half_width - Inches(0.4), Inches(0.5),
                      "\u2713 Onde funciona (auto-cancelamento)",
                      font_size=18, bold=True, color=DARK_GREEN)

    left_lines = [
        ("Teste Chebyshev (n\u22653):", True),
        ("\u03be = 1/n + (\u2016\u03b4\u2016\xb72/d)\u00b2 / (n \xb7 \u03c3\u00b2)", False),
        ("\u03c3\u00b2 tamb\u00e9m usa (2/d)\u00b2", False),
        ("", False),
        ("\u2192 Numerador e denominador cancelam!", True),
        ("O teste principal \u00e9 self-consistent", False),
        ("em QUALQUER dimens\u00e3o", False),
    ]
    tf_left = add_text_box(slide, left_x + Inches(0.3), box_top + Inches(0.7),
                           half_width - Inches(0.6), box_h - Inches(0.9),
                           "", font_size=14)
    first = True
    for text, is_bold in left_lines:
        if first:
            tf_left.paragraphs[0].text = text
            tf_left.paragraphs[0].font.size = Pt(14)
            tf_left.paragraphs[0].font.bold = is_bold
            tf_left.paragraphs[0].font.color.rgb = DARK_GREEN if is_bold else DARK_GRAY
            tf_left.paragraphs[0].font.name = "Calibri"
            first = False
        else:
            c = DARK_GREEN if is_bold else DARK_GRAY
            add_paragraph(tf_left, text, font_size=14, bold=is_bold, color=c,
                          space_before=Pt(4))

    # RIGHT half: red-tinted box
    right_x = Inches(6.8)

    shape = slide.shapes.add_shape(5, right_x, box_top, half_width, box_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
    shape.line.color.rgb = RED
    shape.line.width = Pt(2)

    add_text_box(slide, right_x + Inches(0.2), box_top + Inches(0.1),
                 half_width - Inches(0.4), Inches(0.5),
                 "\u2717 Onde colapsa (3 code paths)",
                 font_size=18, bold=True, color=RED)

    right_lines = [
        ("\u2460 Intersec\u00e7\u00e3o macro-clusters:", True),
        ("   dist(\u03bc\u1d62,\u03bc\u2c7c) vs \u221a\u03c3\u00b2 (escalada)", False),
        ("   \u2192 nunca intersectam", False),
        ("", False),
        ("\u2461 Guard n<3:", True),
        ("   \u03c3\u00b2(escalada) vs r\u2080(fixo)", False),
        ("   \u2192 rejei\u00e7\u00e3o prematura", False),
        ("", False),
        ("\u2462 Life decay:", True),
        ("   \u221a\u03c3\u00b2 e dist em escalas diferentes", False),
        ("   \u2192 morte err\u00e1tica", False),
    ]
    tf_right = add_text_box(slide, right_x + Inches(0.3), box_top + Inches(0.7),
                            half_width - Inches(0.6), box_h - Inches(0.9),
                            "", font_size=13)
    first = True
    for text, is_bold in right_lines:
        if first:
            tf_right.paragraphs[0].text = text
            tf_right.paragraphs[0].font.size = Pt(13)
            tf_right.paragraphs[0].font.bold = is_bold
            tf_right.paragraphs[0].font.color.rgb = RED if is_bold else DARK_GRAY
            tf_right.paragraphs[0].font.name = "Calibri"
            first = False
        else:
            c = RED if is_bold else DARK_GRAY
            add_paragraph(tf_right, text, font_size=13, bold=is_bold, color=c,
                          space_before=Pt(3))

    # Bottom: unified table showing d=2 vs d=17 per code path
    table = add_table(slide, Inches(0.3), Inches(6.1),
                      Inches(12.8), Inches(1.2), 5, 3)

    # Headers
    for i, h in enumerate(["Code Path", "d=2  (2/d)^2=1", "d=17  (2/d)^2=0.014"]):
        table.cell(0, i).text = h
    style_header_row(table, color=GOLD, font_color=BLACK)

    # Data rows
    rows_data = [
        ("Chebyshev (n>=3)", "OK", "OK (auto-cancela)"),
        ("Intersec\u00e7\u00e3o MC", "OK", "FALHA: raio 12% do real"),
        ("Guard n<3", "OK", "FALHA: \u03c3\u00b2 descalibrada vs r\u2080"),
        ("Life decay", "OK", "FALHA: \u221a\u03c3\u00b2 e dist em escalas \u2260"),
    ]
    for r, (path, d2, d17) in enumerate(rows_data, 1):
        set_cell(table, r, 0, path, font_size=12, bold=True,
                 alignment=PP_ALIGN.LEFT)
        set_cell(table, r, 1, d2, font_size=12, color=DARK_GREEN)
        d17_color = DARK_GREEN if "OK" in d17 else RED
        set_cell(table, r, 2, d17, font_size=12, bold=True, color=d17_color)

    # Adjust column widths
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(4.0)
    table.columns[2].width = Inches(5.3)

    # ==================================================================
    # SLIDE 5: AS 5 ADAPTACOES (white bg)
    # ==================================================================
    print("Gerando slide 5: As 5 Adaptacoes...")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "5 Adaptacoes para Alta Dimensao")

    # Table: 6 rows (header + 5 data), 5 columns
    table = add_table(slide, Inches(0.2), Inches(1.5),
                      Inches(12.9), Inches(4.5), 6, 5)

    headers = ["#", "Adaptacao", "Original \u2192 Proposto", "Impacto", "Suporte Teorico"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
    style_header_row(table)

    # Set column widths (approximate via cell width)
    col_widths = [Inches(0.5), Inches(1.8), Inches(5.5), Inches(2.5), Inches(2.6)]
    for i, w in enumerate(col_widths):
        for row_idx in range(6):
            table.cell(row_idx, i).width = w

    data_rows = [
        ["1", "Vari\u00e2ncia",
         "(\u2016\u03b4\u2016\xb72/d)\u00b2 \u2192 Welford dot(\u03b4\u2098\u2091, \u03b4\u2099\u2091\u2097)",
         "Corrige 70x", "Welford 1962, Chan 1983"],
        ["2", "Eccentricidade",
         "(norm\xb72/d)\u00b2 \u2192 \u2016diff\u2016\u00b2",
         "Consistente com \u03c3\u00b2", "Angelov 2014"],
        ["3", "Update",
         "Todos os MC \u2192 S\u00f3 melhor (max \u03c4)",
         "-20% clusters", "Kohonen 1990, NS-TEDA 2024"],
        ["4", "Guard n=1",
         "Sem prote\u00e7\u00e3o \u2192 threshold=13",
         "Seeds sobrevivem", "Reynolds 2009"],
        ["5", "Guard n=2",
         "S\u00f3 \u03c3\u00b2>r\u2080 \u2192 \u03b6>thr AND \u03c3\u00b2\u2265r\u2080",
         "-20% splits", "Reynolds 2009"],
    ]
    for r, row_data in enumerate(data_rows):
        for c, val in enumerate(row_data):
            bold = (c == 0)
            set_cell(table, r + 1, c, val, font_size=12, bold=bold,
                     color=DARK_BLUE if c == 0 else BLACK)

    # Footer
    add_text_box(slide, Inches(0.5), Inches(6.3), Inches(12.5), Inches(0.6),
                 "Cada adaptacao e togglable independentemente (8 variantes V0-V7). "
                 "Ablation study planejado.",
                 font_size=14, italic=True, color=MEDIUM_BLUE)

    # ==================================================================
    # SLIDE 6: PIPELINE STREAMING (white bg)
    # ==================================================================
    print("Gerando slide 6: Pipeline Streaming...")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Pipeline Streaming: Por Que Alta Dimens\u00e3o Importa")

    # ASCII pipeline diagram in monospace text box
    shape = slide.shapes.add_shape(5, Inches(0.5), Inches(1.5),
                                   Inches(12.3), Inches(2.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    shape.line.color.rgb = DARK_BLUE
    shape.line.width = Pt(1.5)

    pipeline_text = (
        "PCAP (CICIoT2023)  \u2192  Kafka  \u2192  FlowConsumer  \u2192  [17 features/flow]  \u2192  MicroTEDAclus  \u2192  Anomalia?\n"
        "                                                              \u2502\n"
        "                                                    WindowAggregator  \u2192  [12-19 feat/janela]  \u2192  MicroTEDAclus  \u2192  Anomalia?"
    )
    add_text_box(slide, Inches(0.7), Inches(1.65), Inches(11.9), Inches(1.7),
                 pipeline_text, font_size=13, bold=True, color=DARK_BLUE,
                 font_name="Courier New")

    # Three key points
    bullets = [
        "17 features por flow \u2192 d=17 \u2192 onde o bug dimensional se manifesta",
        "Prequential (test-then-train): sem data leakage, exige detector O(1)",
        "Gap sem\u00e2ntico (Sommer & Paxson 2010): detector encontra outliers, n\u00e3o ataques",
    ]
    tf = add_text_box(slide, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.5),
                      bullets[0], font_size=15, color=DARK_GRAY)
    for bullet in bullets[1:]:
        add_paragraph(tf, bullet, font_size=15, color=DARK_GRAY, space_before=Pt(10))

    # Footer italic
    add_text_box(slide, Inches(0.5), Inches(6.2), Inches(12.5), Inches(0.6),
                 "Janelas temporais mudam a pergunta: 'este IP est\u00e1 se comportando de forma an\u00f4mala?'",
                 font_size=13, italic=True, color=MEDIUM_BLUE)

    # ==================================================================
    # SLIDE 7: EVIDENCIA (white bg)
    # ==================================================================
    print("Gerando slide 7: Evidencia Experimental...")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Evidencia Experimental")

    # Block 1: C04 small table
    add_text_box(slide, Inches(0.5), Inches(1.4), Inches(5.0), Inches(0.4),
                 "C04: 30 runs \u2014 Proprio vs Original",
                 font_size=15, bold=True, color=DARK_BLUE)

    table = add_table(slide, Inches(0.5), Inches(1.9),
                      Inches(7.5), Inches(1.8), 4, 4)
    for i, h in enumerate(["Config", "Proprio", "Original", "Razao"]):
        table.cell(0, i).text = h
    style_header_row(table)

    c04_data = [
        ["flow r0=0.10", "3.9%", "54.4%", "14x"],
        ["window v1 w=10s", "2.9%", "41.9%", "14x"],
        ["window v1 w=30s", "5.0%", "74.5%", "15x"],
    ]
    for r, row_data in enumerate(c04_data):
        for c, val in enumerate(row_data):
            clr = RED if c == 2 else (DARK_GREEN if c == 1 else BLACK)
            set_cell(table, r + 1, c, val, font_size=12, color=clr,
                     bold=(c == 3))

    # Block 2: Melhor por ataque
    add_text_box(slide, Inches(0.5), Inches(4.0), Inches(12.5), Inches(0.4),
                 "Melhor por ataque (167 exps):",
                 font_size=15, bold=True, color=DARK_BLUE)
    add_text_box(slide, Inches(0.8), Inches(4.45), Inches(12.0), Inches(0.4),
                 "Recon F1=43.7%  |  ICMP Recall=50%  |  SYN Recall=54% (FPR=33%)  |  TCP = 0%",
                 font_size=14, color=DARK_GRAY)

    # Block 3: Gap semantico
    shape = slide.shapes.add_shape(5, Inches(0.5), Inches(5.1),
                                   Inches(12.3), Inches(1.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE1)
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(1.5)

    tf = add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11.8), Inches(0.4),
                      "Gap semantico nos dados:",
                      font_size=14, bold=True, color=GOLD)
    gap_lines = [
        "Per-flow: anomaly rate ~3.5% COM ou SEM ataque (C01)",
        "Janelas: SYN 3%\u219254%, Recon 4%\u219245% \u2014 mas FPR sobe (14% em w=10s v2)",
        "Gap: sem baseline nao-supervisionado direto no CICIoT2023",
    ]
    for line in gap_lines:
        add_paragraph(tf, line, font_size=12, color=DARK_GRAY, space_before=Pt(4))

    # ==================================================================
    # SLIDE 8: RESULTADOS NOVOS (white bg)
    # ==================================================================
    print("Gerando slide 8: Resultados Novos...")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Resultados Novos (Exp 1 + Exp 2)")

    # Left: Exp 1 summary
    shape = slide.shapes.add_shape(5, Inches(0.3), Inches(1.5),
                                   Inches(6.3), Inches(3.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(1.5)

    add_text_box(slide, Inches(0.5), Inches(1.55), Inches(5.9), Inches(0.4),
                 "Exp 1: Sweep Dimensional (1440 runs)",
                 font_size=16, bold=True, color=DARK_GREEN)

    exp1_lines = [
        "Com r\u2080\u22651.0 (guard neutralizado):",
        "  V0 FPR = 0% em QUALQUER d",
        "  \u2192 Chebyshev auto-cancela (CONFIRMADO)",
        "",
        "Com r\u2080=0.001 (mal calibrado):",
        "  V0 FPR = 72-99.8% em todo d",
        "  \u2192 Guard n<3 domina tudo",
        "",
        "V7 est\u00e1vel ~0.1% em TODOS os r\u2080 e d",
    ]
    tf = add_text_box(slide, Inches(0.5), Inches(2.1), Inches(5.9), Inches(2.5),
                      exp1_lines[0], font_size=13, color=DARK_GRAY)
    for line in exp1_lines[1:]:
        bold = "\u2192" in line or "CONFIRMADO" in line
        add_paragraph(tf, line, font_size=13, bold=bold,
                      color=DARK_GREEN if bold else DARK_GRAY, space_before=Pt(3))

    # Right: Exp 2 summary
    shape = slide.shapes.add_shape(5, Inches(6.8), Inches(1.5),
                                   Inches(6.3), Inches(3.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
    shape.line.color.rgb = RED
    shape.line.width = Pt(1.5)

    add_text_box(slide, Inches(7.0), Inches(1.55), Inches(5.9), Inches(0.4),
                 "Exp 2: Abla\u00e7\u00e3o V0-V7 (240 runs, d=17)",
                 font_size=16, bold=True, color=RED)

    exp2_lines = [
        "V7 (full):     FPR = 0.1%  \u2713",
        "V0 (original): FPR = 0%    \u2713 (r\u2080=1.0)",
        "V5 (guard n1): FPR = 0.4%  \u2713",
        "",
        "V2 (ecc only): FPR = 49.7% \u2717",
        "V1 (Welford):  FPR = 98.9% \u2717\u2717",
        "V3 (Welf+ecc): FPR = 99.8% \u2717\u2717\u2717",
        "",
        "Welford SOZINHO piora!",
        "Adapta\u00e7\u00f5es s\u00e3o ACOPLADAS",
    ]
    tf = add_text_box(slide, Inches(7.0), Inches(2.1), Inches(5.9), Inches(2.5),
                      exp2_lines[0], font_size=13, color=DARK_GRAY)
    for line in exp2_lines[1:]:
        bold = "SOZINHO" in line or "ACOPLADAS" in line
        c = RED if ("\u2717" in line or bold) else DARK_GREEN if "\u2713" in line else DARK_GRAY
        add_paragraph(tf, line, font_size=13, bold=bold, color=c, space_before=Pt(3))

    # Bottom: status + stats
    shape = slide.shapes.add_shape(5, Inches(0.3), Inches(5.2),
                                   Inches(12.8), Inches(1.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    shape.line.color.rgb = DARK_GRAY
    shape.line.width = Pt(1)

    tf = add_text_box(slide, Inches(0.5), Inches(5.3), Inches(12.4), Inches(0.4),
                      "Estat\u00edstica: Friedman \u03c7\u00b2=203.2, p<10\u207b\u2074\u2070 | ANOVA F=19381, p=0 | CD=1.917 (Nem\u00e9nyi)",
                      font_size=12, bold=True, color=DARK_BLUE)
    add_paragraph(tf,
                  "Exp 3: IF e OC-SVM parciais (12/42). Adicionando Half-Space Trees (streaming genu\u00edno) \u2014 rodando no Linux",
                  font_size=13, bold=True, color=ORANGE, space_before=Pt(8))
    add_paragraph(tf,
                  "Achado principal: as adapta\u00e7\u00f5es formam um sistema acoplado. "
                  "Ligar uma isoladamente pode PIORAR o resultado.",
                  font_size=12, italic=True, color=DARK_GRAY, space_before=Pt(6))

    # ==================================================================
    # SLIDE 9: PROXIMOS PASSOS E DECISOES (dark blue bg)
    # ==================================================================
    print("Gerando slide 9: Proximos Passos...")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BLUE)

    add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(1.0),
                 "Pr\u00f3ximos Passos e Decis\u00f5es",
                 font_size=40, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    questions = [
        "1. Auto-cancelamento + acoplamento + estudo emp\u00edrico sustentam um paper SoftCom?",
        "2. Foco: an\u00e1lise dimensional pura OU incluir resultados de detec\u00e7\u00e3o IoT (F1=43.7%)?",
        "3. Timeline: paper SoftCom para quando? Disserta\u00e7\u00e3o para agosto?",
    ]

    for i, question in enumerate(questions):
        y = Inches(2.2) + Inches(i * 1.5)

        shape = slide.shapes.add_shape(5, Inches(1.5), y,
                                       Inches(10.3), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = MEDIUM_BLUE
        shape.line.fill.background()

        add_text_box(slide, Inches(1.8), y + Inches(0.15),
                     Inches(9.8), Inches(0.8),
                     question, font_size=19, bold=True, color=WHITE)

    add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.6),
                 "Obrigado. Abro para discuss\u00e3o.",
                 font_size=18, italic=True, color=LIGHT_BLUE,
                 alignment=PP_ALIGN.CENTER)

    # Save
    prs.save(str(OUTPUT_PATH))
    print(f"\nApresentacao salva em: {OUTPUT_PATH}")
    print(f"Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
